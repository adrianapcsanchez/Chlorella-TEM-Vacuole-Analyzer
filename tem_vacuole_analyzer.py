"""
====================================================================
 TEM VACUOLE AREA ANALYZER - by Adriana Pereira C. Sanchez
====================================================================

What this script does:
 1. Asks you to select the folder containing the images (.tif/.tiff/.png/.jpg).
 2. Displays a window where you can REVIEW/ADJUST the
    segmentation parameters (you do not need to edit the code; see the "PARAMETERS" section
    below).
 3. For each image:
      a) Calibrates image brightness/contrast to a common standard (so brighter and darker images from the same dataset are treated more consistently).
      b) Reads the scale value (ex.: "1 um", "500 nm") in the lower-left corner using OCR and calculates spatial calibration from the scale bar 
      c) Segments the cell(s) in the image (up to 1 to 10 cells).
      d) Opens a review window to add, draw, delete, or split cell contours.
      e) Runs automatic vacuole detection, then opens every image for review.
         A adds from one click, B draws a boundary, and C erases a contour.
      f) Calculates each cell area, each vacuole area, the sum
         of vacuole areas, and the vacuole/cell ratio (%).
      g) Creates a verification image with a GREEN cell outline and
         RED outlines around detected vacuoles.
 4. At the end, it generates:
      - An Excel spreadsheet (.xlsx) with all measurements (one row per detected cell).
      - A PowerPoint presentation (.pptx) with two annotated images per
        slide (side by side), each with the file name above and the
        data table below each image for visual verification.

------------------------------------------------------------------
 IMPORTANT - ABOUT AUTOMATIC RESULTS
------------------------------------------------------------------
Automatic TEM image segmentation is never 100% perfect. Therefore,
the script always generates annotated images (green/red) and a PowerPoint report. Use them to visually verify every result before using the measurements in a publication.

------------------------------------------------------------------
 INSTALLATION (run once in a terminal/command prompt)
------------------------------------------------------------------
pip install numpy opencv-python tifffile imagecodecs pytesseract pandas openpyxl python-pptx pillow matplotlib

In addition, pytesseract requires the Tesseract-OCR application to be installed
on the computer (it is not only a Python library):
  - Windows: download the installer from
    https://github.com/UB-Mannheim/tesseract/wiki  and install it.
    If the script does not find Tesseract automatically, set the
    path manually in CONFIG below (TESSERACT_CMD).
  - Mac: brew install tesseract
  - Linux: sudo apt install tesseract-ocr

If you do not install Tesseract, the script still works. It will only
ask you to ENTER the scale value manually whenever OCR cannot read it in
an image.
====================================================================
"""

__version__ = "1.0.0"
__author__ = "Adriana Augusta Pereira Contreras Sanchez"

import os
import re
import sys
import traceback
import shutil
from pathlib import Path

import numpy as np
import cv2

# ---------------------------------------------------------------
# PARAMETERS - default values (can be adjusted from the program window without editing this file; see edit_parameters_gui)
# ---------------------------------------------------------------
class CONFIG:
    # Path to the Tesseract executable (leave None for automatic detection).
    # Windows example: r"C:\Program Files\Tesseract-OCR\tesseract.exe"
    TESSERACT_CMD = None

    # --- brightness/contrast normalization across datasets ---
    # Each image batch may have a different overall brightness -
    # some images may be brighter and others darker. This normalization
    # stretches the histogram of EACH image to a common reference range
    # BEFORE segmentation so the same parameters behave more consistently
    # across datasets.
    NORMALIZE_BRIGHTNESS = True
    BRIGHTNESS_LOW_PERCENTILE = 1     # LOWER = more sensitive to isolated very dark noise pixels
    BRIGHTNESS_HIGH_PERCENTILE = 99     # HIGHER = more sensitive to isolated very bright noise pixels
    BRIGHTNESS_LOW_REFERENCE = 8           # to where the "low" and stretched percentile (0–255)
    BRIGHTNESS_HIGH_REFERENCE = 245          # to where the "high" and stretched percentile (0–255)
    AUTO_GAMMA = True              # automatically brightens genuinely dark images
    AUTO_GAMMA_TARGET_MEDIAN = 145 # median-target after normalization
    AUTO_GAMMA_TRIGGER_MEDIAN = 125# apply gamma only if the normalized image remains dark
    AUTO_GAMMA_MIN = 0.68          # limit that prevents excessive brightening

    # --- cell segmentation ---
    CLAHE_CLIP = 2.0              # HIGHER = stronger local contrast enhancement (may increase noise)
    CLAHE_TILE = 16               #local contrast enhancement block size
    GAUSS_BLUR_CELL = 11           # HIGHER = stronger smoothing before separating cell/background
    DENOISE_MIN_COMPONENT_FRAC = 0.0003  # HIGHER = removes larger isolated components before cell assembly
    CLOSE_KERNEL_CELL = 41       # HIGHER = joins more distant pieces of the SAME cell (but may
                                    # also merge two nearby cells if set too high)
    OPEN_KERNEL_CELL = 11          # HIGHER = removes thicker thin bridges/tendrils connecting the
                                    # cell to the background (but may remove legitimate thin regions)
    MIN_CELL_AREA_FRAC = 0.01     # LOWER = accepts smaller cells (fraction of image area)
    MAX_CELL_AREA_FRAC = 0.85     # HIGHER = accepts larger cells (up to this fraction of image area)
    MIN_CELL_SOLIDITY = 0.75      # HIGHER = requires a rounder/more filled cell (rejects
                                    # risks and folds, but can reject cells with
                                    # legitimate indentations)
    MAX_CELL_ASPECT_RATIO = 2.5   # LOWER = requires rounder cells (rejects elongated shapes)
    MAX_CELLS_PER_IMAGE = 4       # maximum number of accepted cells per image (largest first)
    CELL_ROUND_HULL = False       # smooths small indentations in the cell boundary
    CELL_HULL_MAX_EXPANSION = 0.12  # prevents the hull from increasing cell area by more than 12%
    CELL_SMOOTH_KERNEL = 11      # final smoothing of the cell boundary
    # Outer-contour refinement: closes inward indentations caused
    # by bright vacuoles near the membrane without converting the cell into a convex hull.
    # Outer-boundary closing at a FIXED scale.
    # 51 px closes large indentations without converting the cell into a convex hull.
    CELL_REENTRANCE_CLOSE_PX = 56

    CELL_OUTER_MAX_EXPANSION = 0.20
    CELL_EDGE_SMOOTH_PX = 7

    # --- vacuole segmentation (single detector, contour red) ---
    # Geometry and the primary criterion. Intensity/contrast come into play later
    # only to help distinguish between the vacuole and the cytoplasm.
    VACUOLE_THRESH_OFFSET = 0
    VACUOLE_CLOSE_KERNEL = 15
    # 5% allows recovery of some smaller real vacuoles. Between 5 and 8%, the
    # candidate must pass much stricter shape/intensity filters.
    MIN_VACUOLE_AREA_FRAC = 0.05
    VACUOLE_STRICT_SMALL_BELOW = 0.08
    MAX_VACUOLE_AREA_FRAC = 0.90
    MIN_VACUOLE_CIRCULARITY = 0.24
    MIN_VACUOLE_SOLIDITY = 0.72
    MIN_VACUOLE_SHAPE_SCORE = 0.47
    MAX_VACUOLE_ASPECT_RATIO = 3.2
    MAX_VACUOLES_PER_CELL = 5

    # --- manual seed-guided vacuole segmentation (V33 mode A) ---
    # The user marks one point inside each vacuole. The program then grows a
    # connected region with a similar gray tone, always restricted to the cell.
    MANUAL_SEED_PATCH_RADIUS = 5
    MANUAL_SEED_TOLERANCES = (5, 7, 9, 12, 16, 21, 27, 34)
    MANUAL_SEED_MIN_AREA_FRAC = 0.004
    MANUAL_SEED_MAX_AREA_FRAC = 0.90
    MANUAL_SEED_MIN_CIRCULARITY = 0.18
    MANUAL_SEED_MIN_SOLIDITY = 0.68
    MANUAL_SEED_MAX_ASPECT_RATIO = 3.5
    MANUAL_SEED_SMOOTH_FRAC = 0.018

    # In this version, a vacuole is treated as a relatively light-gray region inside the cell.
    # Intensity does not create a candidate by itself; it confirms the geometry.
    VACUOLE_CONTRAST_RING_FRAC = 0.035
    VACUOLE_MIN_BRIGHTNESS_DELTA = 4.0   # median(vacuole) - median(surrounding cytoplasm)
    VACUOLE_MIN_SUPPORT_RATIO = 0.50     # fraction of the interior supported by bright pixels before closing
    VACUOLE_MAX_ROBUST_SPREAD = 130.0   # maximum p90-p10 spread; prevents inclusion of heterogeneous cytoplasm
    VACUOLE_MIN_CELL_PERCENTILE = 40.0  # candidate median must be >= cell p40

    # small vacuoles (5-8%) They only get in if they are particularly convincing.
    VACUOLE_SMALL_MIN_CIRCULARITY = 0.38
    VACUOLE_SMALL_MIN_SOLIDITY = 0.80
    VACUOLE_SMALL_MAX_ASPECT_RATIO = 2.8
    VACUOLE_SMALL_MIN_SUPPORT_RATIO = 0.65
    VACUOLE_SMALL_MAX_ROBUST_SPREAD = 100.0

    # Complementary search using LOCAL contrast (always white regions).
    VACUOLE_LOCAL_SIGMA_FRAC_1 = 0.06
    VACUOLE_LOCAL_SIGMA_FRAC_2 = 0.14

    # A vacuole must remain inside the cell; small internal margin for smaller objects.
    VACUOLE_BORDER_CONTACT_MAX = 0.00
    VACUOLE_HUGE_BORDER_CONTACT_MAX = 0.00
    VACUOLE_MIN_BORDER_MARGIN_FRAC = 0.02

    # --- Grayscale classification inside the cell ---
    #   dark = cytoplasm / dense material
    #   intermediate gray = vacuole
    #   very bright = lipids / white inclusions
    # Thresholds are recalculated for every cell using 1D k-means.
    VACUOLE_TONE_K = 3
    VACUOLE_TONE_MIN_GAP = 8.0
    VACUOLE_GRAY_LOWER_PAD = 3.0
    VACUOLE_GRAY_UPPER_PAD = 2.0
    VACUOLE_MAX_WHITE_FRACTION = 0.12
    VACUOLE_MAX_DARK_FRACTION = 0.28
    VACUOLE_GRAY_MIN_SUPPORT = 0.62
    VACUOLE_GRAY_SMALL_MIN_SUPPORT = 0.72
    VACUOLE_GRAY_MAX_SPREAD = 95.0
    VACUOLE_GRAY_SMALL_MAX_SPREAD = 75.0
    VACUOLE_GRAY_CLOSE_KERNEL = 7
    VACUOLE_GRAY_OPEN_KERNEL = 5

    # adaptive tone + texture calibration (adaptive per cell)
    # The vacuole is expected to be intermediate gray and relatively smooth/homogeneous.
    # Cytoplasm is usually darker/more textured; lipids are very bright.
    VACUOLE_TONE_LOW_PERCENTILE = 15.0
    VACUOLE_TONE_HIGH_PERCENTILE = 88.0
    VACUOLE_TEXTURE_WINDOW_FRAC = 0.018
    VACUOLE_TEXTURE_PERCENTILE = 45.0
    VACUOLE_MIN_ELLIPSE_IOU = 0.72
    VACUOLE_ELLIPSE_MIN_SOLIDITY = 0.80
    VACUOLE_ELLIPSE_MAX_EXPANSION = 0.15

    # --- complementary SHAPE + relative-brightness detector ---
    SHAPE_DETECTOR_ENABLED = True
    SHAPE_MIN_AREA_FRAC = 0.02
    SHAPE_STRICT_BELOW_FRAC = 0.08
    SHAPE_MAX_AREA_FRAC = 0.90
    SHAPE_MIN_CIRCULARITY = 0.24
    SHAPE_MIN_SOLIDITY = 0.72
    SHAPE_MAX_ASPECT_RATIO = 3.4
    SHAPE_SMALL_MIN_CIRCULARITY = 0.34
    SHAPE_SMALL_MIN_SOLIDITY = 0.78
    SHAPE_SMALL_MAX_ASPECT_RATIO = 2.8
    SHAPE_MIN_LOCAL_BRIGHTNESS_DELTA = 1.5
    SHAPE_RING_FRAC = 0.035
    SHAPE_THRESHOLD_PERCENTILES = (32, 42, 52, 62, 72)
    SHAPE_MORPH_CLOSE_PX = 7
    SHAPE_MORPH_OPEN_PX = 3

    # --- detector fusion ---
    HYBRID_USE_V24 = True
    HYBRID_DUPLICATE_OVERLAP = 0.86
    HYBRID_DUPLICATE_AREA_RATIO_MIN = 0.70
    HYBRID_DUPLICATE_AREA_RATIO_MAX = 1.43
    REFINE_CONTOUR_WITH_EDGES = False
    REFINE_CONTOUR_ITERATIONS = 18
    REFINE_CONTOUR_MARGIN_FRAC = 0.035
    REFINE_CONTOUR_MAX_AREA_CHANGE = 0.20
    REFINE_CONTOUR_MIN_SOLIDITY = 0.68

    # --- scale (bar + label in the lower-left corner) ---
    # reference is available as a fallback when automatic measurement fails
    # for the exact bar length). The script reads the
    # LABEL (e.g., "1 um", "500 nm") by OCR in the lower-left corner.
    SCALEBAR_LENGTH_PX = 500   # fallback ONLY if the bar cannot be measured
    SCALEBAR_SEARCH_WIDTH_FRAC = 0.50
    SCALEBAR_SEARCH_HEIGHT_FRAC = 0.30
    SCALEBAR_DETECT_LENGTH = True
    SCALEBAR_MIN_WIDTH_FRAC = 0.025
    SCALEBAR_MAX_HEIGHT_FRAC = 0.025
    SCALEBAR_MIN_ASPECT_RATIO = 5.0
    SCALEBAR_TEXT_HEIGHT_FRAC = 0.10

    # When OCR fails, manual calibration is available for that image.
    CALIBRATION_FAILURE_MODE = "manual_required"

    # --- Output ---
    DPI_FIGURA = 150




# ---------------------------------------------------------------
# Optional/heavier dependencies
# ---------------------------------------------------------------
try:
    import tifffile
except ImportError:
    print("ERROR: missing dependency 'tifffile'. Run: pip install tifffile imagecodecs")
    sys.exit(1)

try:
    import pytesseract

    def _autodetect_tesseract():
        candidates = []
        if CONFIG.TESSERACT_CMD:
            candidates.append(CONFIG.TESSERACT_CMD)

        detected_path = shutil.which("tesseract")
        if detected_path:
            candidates.append(detected_path)

        if os.name == "nt":
            localapp = os.environ.get("LOCALAPPDATA", "")
            programfiles = os.environ.get("ProgramFiles", r"C:\\Program Files")
            programfiles86 = os.environ.get("ProgramFiles(x86)", r"C:\\Program Files (x86)")
            candidates.extend([
                os.path.join(localapp, "Tesseract-OCR", "tesseract.exe"),
                os.path.join(programfiles, "Tesseract-OCR", "tesseract.exe"),
                os.path.join(programfiles86, "Tesseract-OCR", "tesseract.exe"),
            ])

        for c in candidates:
            if c and Path(c).exists():
                pytesseract.pytesseract.tesseract_cmd = str(c)
                return str(c)
        return None

    _TESSERACT_FOUND = _autodetect_tesseract()
    TESSERACT_OK = bool(_TESSERACT_FOUND)
except ImportError:
    _TESSERACT_FOUND = None
    TESSERACT_OK = False

import pandas as pd
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt

import tkinter as tk
from tkinter import filedialog, simpledialog, messagebox, ttk

import matplotlib
try:
    import matplotlib.pyplot as plt
except Exception:
    # Fallback without forcing a backend. On Windows, the Tk backend is
    # normally selected automatically.
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt


# =================================================================
# 1) IMAGE LOADING (supports .tif/.tiff with JEOL/Gatan metadata,
#    as well as standard PNG/JPG files)
# =================================================================
def load_grayscale_image(path):
    """Returns the image as a 2D NumPy array (uint8, grayscale)."""
    ext = Path(path).suffix.lower()
    if ext in (".tif", ".tiff"):
        try:
            image = tifffile.imread(path)
        except Exception:
            # Useful fallback for LZW TIFF files when imagecodecs is unavailable.
            image = np.array(Image.open(path).convert("L"))
    else:
        image = np.array(Image.open(path).convert("L"))

    # If the image has more than 2 dimensions (RGB or stack), reduce it to 2D grayscale.
    if image.ndim == 3:
        if image.shape[-1] in (3, 4):
            image = cv2.cvtColor(image[..., :3], cv2.COLOR_RGB2GRAY)
        else:
            image = image[0]

    # Normalize to uint8 if the source is 16-bit or floating point.
    if image.dtype != np.uint8:
        image = cv2.normalize(image, None, 0, 255, cv2.NORM_MINMAX).astype(np.uint8)

    return image


def normalize_image_brightness(image):
    """Normalize brightness/contrast BEFORE segmentation.

    1) Stretch the histogram using the configured percentiles.
    2) If the image is still dark, apply adaptive gamma brightening.
    3) Local CLAHE is subsequently applied within the segmentation process.

    The original image is NOT overwritten; the normalized image is used only
    for analysis/segmentation. In grayscale TEM images there is no
    color saturation in the usual sense; the useful equivalent is controlled
    contrast and luminance enhancement.
    """
    if not CONFIG.NORMALIZE_BRIGHTNESS:
        return image.copy()

    low = float(np.percentile(image, CONFIG.BRIGHTNESS_LOW_PERCENTILE))
    high = float(np.percentile(image, CONFIG.BRIGHTNESS_HIGH_PERCENTILE))
    if high - low < 1:
        return image.copy()

    normalized = (image.astype(np.float32) - low) * (
        CONFIG.BRIGHTNESS_HIGH_REFERENCE - CONFIG.BRIGHTNESS_LOW_REFERENCE
    ) / (high - low) + CONFIG.BRIGHTNESS_LOW_REFERENCE
    normalized = np.clip(normalized, 0, 255).astype(np.uint8)

    if CONFIG.AUTO_GAMMA:
        median = float(np.median(normalized))
        if 1.0 < median < CONFIG.AUTO_GAMMA_TRIGGER_MEDIAN:
            target = float(CONFIG.AUTO_GAMMA_TARGET_MEDIAN) / 255.0
            current = median / 255.0
            gamma_value = np.log(target) / np.log(current)
            gamma_value = float(np.clip(gamma_value, CONFIG.AUTO_GAMMA_MIN, 1.0))
            table = np.array([
                np.clip(((i / 255.0) ** gamma_value) * 255.0, 0, 255)
                for i in range(256)
            ], dtype=np.uint8)
            normalized = cv2.LUT(normalized, table)

    return normalized


# =================================================================
# 2) SCALE: read the label in the lower-left corner (OCR), and use the scale-bar length (with automatic bar detection when possible)
# =================================================================
def _normalize_scale_text(text):
    """Normalizes common OCR errors in the scale label."""
    if text is None:
        return ""
    t = text.lower().replace("μ", "µ")
    t = t.replace("|", "1")
    t = re.sub(r"\s+", " ", t)
    # common OCR errors: 'n m' -> 'nm', 'u m' -> 'um'
    t = re.sub(r"\bn\s*m\b", "nm", t)
    t = re.sub(r"\b[uµ]\s*m\b", "um", t)
    return t.strip()


def _extract_scale_value_unit(text):
    """Extracts e.g. 500 nm, 0.5 um, 1 µm. Accepts minor OCR errors."""
    t = _normalize_scale_text(text)
    m = re.search(r"([0-9]+(?:[.,][0-9]+)?)\s*(nm|um|µm)", t, re.IGNORECASE)
    if not m:
        return None
    value = float(m.group(1).replace(",", "."))
    read_unit = m.group(2).lower()
    unit = "nm" if read_unit == "nm" else "um"
    # Reject obviously implausible OCR readings.
    if value <= 0 or value > 100000:
        return None
    return value, unit


def _ocr_scale_text_multi(images):
    """Runs OCR on several preprocessed versions of the black scale label."""
    if not TESSERACT_OK:
        return None
    if not isinstance(images, (list, tuple)):
        images = [images]

    ocr_configs = [
        "--psm 7",
        "--psm 6",
        "--psm 11",
        "--psm 13",
    ]
    whitelist = "0123456789.,nmuµ "
    for img in images:
        for cfg in ocr_configs:
            try:
                t = pytesseract.image_to_string(
                    img,
                    config=f"{cfg} -c tessedit_char_whitelist={whitelist}",
                    timeout=4,
                ).strip()
            except Exception:
                continue
            r = _extract_scale_value_unit(t)
            if r:
                return r
    return None


def _detect_scale_bar(image):
    """Detect the BLACK horizontal scale bar in the lower-left corner.

    Return a dictionary with x, y, width, and height in image coordinates, or None.
    """
    h, w = image.shape
    rx = max(50, int(w * CONFIG.SCALEBAR_SEARCH_WIDTH_FRAC))
    ry = max(50, int(h * CONFIG.SCALEBAR_SEARCH_HEIGHT_FRAC))
    y0 = h - ry
    roi = image[y0:h, 0:rx]

    # The scale bar is black on a light background: inverted Otsu highlights dark elements.
    smoothed = cv2.GaussianBlur(roi, (3, 3), 0)
    _, bw = cv2.threshold(smoothed, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)

    # Remove small dots/text and preserve long horizontal segments.
    kw = max(21, int(rx * CONFIG.SCALEBAR_MIN_WIDTH_FRAC))
    kh = 3
    horiz = cv2.morphologyEx(
        bw,
        cv2.MORPH_OPEN,
        cv2.getStructuringElement(cv2.MORPH_RECT, (kw, kh)),
    )

    contours, _ = cv2.findContours(horiz, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    candidates = []
    for c in contours:
        x, y, ww, hh = cv2.boundingRect(c)
        asp = ww / max(hh, 1)
        if ww < w * CONFIG.SCALEBAR_MIN_WIDTH_FRAC:
            continue
        if hh > h * CONFIG.SCALEBAR_MAX_HEIGHT_FRAC:
            continue
        if asp < CONFIG.SCALEBAR_MIN_ASPECT_RATIO:
            continue
        # Favor bars located in the lower part of the ROI.
        score = ww * (1.0 + 0.35 * ((y + hh) / max(ry, 1)))
        candidates.append((score, x, y + y0, ww, hh))

    if not candidates:
        return None
    candidates.sort(reverse=True)
    _, x, y, ww, hh = candidates[0]
    return {"x": int(x), "y": int(y), "w": int(ww), "h": int(hh)}


def _preprocess_scale_text(image, scale_bar=None):
    """Create several preprocessed versions of the black scale text for OCR."""
    h, w = image.shape

    if scale_bar is not None:
        bx, by, bw, bh = scale_bar["x"], scale_bar["y"], scale_bar["w"], scale_bar["h"]
        # The label is usually immediately ABOVE the scale bar.
        x_margin = max(15, int(0.18 * bw))
        x0 = max(0, bx - x_margin)
        x1 = min(w, bx + int(1.8 * bw))
        height = max(int(h * CONFIG.SCALEBAR_TEXT_HEIGHT_FRAC), int(5 * max(bh, 1)))
        y0 = max(0, by - height)
        y1 = min(h, by + max(3, bh))
        crop = image[y0:y1, x0:x1]
    else:
        ry = int(h * CONFIG.SCALEBAR_SEARCH_HEIGHT_FRAC)
        rx = int(w * CONFIG.SCALEBAR_SEARCH_WIDTH_FRAC)
        crop = image[h - ry:h, 0:rx]

    if crop.size == 0:
        return []

    # The text crop does not need extreme enlargement: 2.5-3x is enough
    largest_dim = max(crop.shape)
    scale_factor = max(2.0, min(4.0, 1200.0 / max(largest_dim, 1)))
    big = cv2.resize(crop, None, fx=scale_factor, fy=scale_factor, interpolation=cv2.INTER_CUBIC)

    # The label is black; preserve dark text and test several binarizations.
    big = cv2.normalize(big, None, 0, 255, cv2.NORM_MINMAX)
    variants = [big]

    _, otsu = cv2.threshold(big, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    variants.append(otsu)

    # Adaptive thresholding helps when the corner background is not uniform.
    block_size = 51 if min(big.shape) >= 51 else max(11, (min(big.shape)//2)*2 - 1)
    if block_size >= 11 and block_size % 2 == 1:
        adap = cv2.adaptiveThreshold(
            big, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
            cv2.THRESH_BINARY, block_size, 11
        )
        variants.append(adap)

    # A minimal dilation makes thin black characters easier to read.
    inv = cv2.bitwise_not(otsu)
    inv = cv2.dilate(inv, np.ones((2, 2), np.uint8), iterations=1)
    variants.append(cv2.bitwise_not(inv))

    return variants


def read_scale_label(image):
    """Robust OCR for black scale text in the lower-left corner."""
    if not TESSERACT_OK:
        return None

    scale_bar = _detect_scale_bar(image)

    if scale_bar is not None:
        reading = _ocr_scale_text_multi(_preprocess_scale_text(image, scale_bar))
        if reading:
            return reading

    h, w = image.shape
    rx = int(w * max(CONFIG.SCALEBAR_SEARCH_WIDTH_FRAC, 0.55))

    for frac_h in (0.12, 0.18, 0.25, 0.32):
        ry = max(40, int(h * frac_h))
        crop = image[h-ry:h, 0:rx]
        if crop.size == 0:
            continue

        scale_factor = max(2.0, min(4.0, 1400.0 / max(crop.shape)))
        big = cv2.resize(crop, None, fx=scale_factor, fy=scale_factor, interpolation=cv2.INTER_CUBIC)
        big = cv2.normalize(big, None, 0, 255, cv2.NORM_MINMAX)

        variants = [big]
        _, otsu = cv2.threshold(big, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        variants.extend([otsu, cv2.bitwise_not(otsu)])

        block = 51 if min(big.shape) >= 51 else 31
        if block % 2 == 0:
            block += 1
        adap = cv2.adaptiveThreshold(
            big, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
            cv2.THRESH_BINARY, block, 9
        )
        variants.extend([adap, cv2.bitwise_not(adap)])

        reading = _ocr_scale_text_multi(variants)
        if reading:
            return reading

    return None




def automatic_scale_calibration(image):
    """Read the scale value/unit and, when possible, measure the black bar length in pixels."""
    reading = read_scale_label(image)
    if reading is None:
        return None

    value, unit = reading
    value_um = value if unit == "um" else value / 1000.0

    scale_bar = _detect_scale_bar(image) if CONFIG.SCALEBAR_DETECT_LENGTH else None
    if scale_bar is not None and scale_bar["w"] >= 10:
        bar_length_px = float(scale_bar["w"])
        bar_method = "black bar detected"
    else:
        bar_length_px = float(CONFIG.SCALEBAR_LENGTH_PX)
        bar_method = "fixed-length fallback"

    um_per_px = value_um / bar_length_px
    return {
        "um_per_px": um_per_px,
        "method": "automatic",
        "details": (
            f"read='{value} {unit}', {bar_method}={bar_length_px:.0f}px"
        ),
    }


def manual_scale_calibration(image, file_name):
    """Manual fallback: automatic bar plus typed value, or two clicks if bar detection fails."""
    h, w = image.shape
    scale_bar = _detect_scale_bar(image)

    ry = int(h * max(CONFIG.SCALEBAR_SEARCH_HEIGHT_FRAC, 0.30))
    rx = int(w * max(CONFIG.SCALEBAR_SEARCH_WIDTH_FRAC, 0.55))
    crop = image[h-ry:h, 0:rx]

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 6))
    ax1.imshow(image, cmap="gray")
    ax2.imshow(crop, cmap="gray")
    ax2.set_title("Lower-left corner zoom")
    ax1.set_title(file_name)

    if scale_bar is not None:
        ax1.add_patch(
            plt.Rectangle(
                (scale_bar["x"], scale_bar["y"]),
                scale_bar["w"], max(scale_bar["h"], 3),
                linewidth=2, edgecolor="red", facecolor="none"
            )
        )

    fig.tight_layout()
    plt.show(block=False)
    plt.pause(0.2)

    if scale_bar is not None and scale_bar["w"] >= 10:
        text = simpledialog.askstring(
            "Manual calibration required",
            f"OCR failed for:\\n{file_name}\\n\\n"
            f"Black scale bar detected: {scale_bar['w']} px.\\n"
            "Enter the printed value, for example: 0.5 um, 1 um, or 500 nm."
        )
        plt.close(fig)

        if not text:
            return None

        r = _extract_scale_value_unit(text)
        if r is None:
            return None

        value, unit = r
        value_um = value if unit == "um" else value / 1000.0
        return {
            "um_per_px": value_um / float(scale_bar["w"]),
            "method": "manual_text_auto_bar",
            "details": f"{value} {unit}; detected bar={scale_bar['w']}px",
        }

    ax1.set_title(f"{file_name}\\nClick the TWO endpoints of the scale bar")
    plt.show(block=False)
    points = plt.ginput(2, timeout=0)
    plt.close(fig)

    if len(points) != 2:
        return None

    (x1, y1), (x2, y2) = points
    distance_px = float(np.hypot(x2-x1, y2-y1))
    if distance_px < 2:
        return None

    text = simpledialog.askstring(
        "Scale value",
        f"Marked length: {distance_px:.1f} px.\\n"
        "Enter the actual value: 0.5 um, 1 um, 500 nm, etc."
    )
    if not text:
        return None

    r = _extract_scale_value_unit(text)
    if r is None:
        return None

    value, unit = r
    value_um = value if unit == "um" else value / 1000.0
    return {
        "um_per_px": value_um / distance_px,
        "method": "manual_clicks",
        "details": f"{distance_px:.1f}px = {value} {unit}",
    }


# =================================================================
#
# =================================================================
# 3) CELL AND VACUOLE SEGMENTATION
# =================================================================
def _odd_kernel(value, minimum=3):
    """Ensure an odd morphological kernel greater than or equal to the minimum."""
    k = int(round(value))
    k = max(minimum, k)
    if k % 2 == 0:
        k += 1
    return k


def _contour_metrics(c):
    area = cv2.contourArea(c)
    perimeter = cv2.arcLength(c, True)
    circularity = (4 * np.pi * area / (perimeter ** 2)) if perimeter > 0 else 0.0
    hull = cv2.convexHull(c)
    hull_area = cv2.contourArea(hull)
    solidity = (area / hull_area) if hull_area > 0 else 0.0
    (_, _), (rw, rh), _ = cv2.minAreaRect(c)
    max_length, minor_axis = max(rw, rh), max(min(rw, rh), 1e-6)
    aspect_ratio = max_length / minor_axis
    return area, circularity, solidity, aspect_ratio


def _accept_vacuole_candidate(c, cell_area, parameters):
    """Applies area and shape filters to a vacuole contour."""
    area, circularity, solidity, aspect_ratio = _contour_metrics(c)
    area_fraction = area / max(cell_area, 1e-9)
    if not (parameters["min_area"] <= area_fraction <= parameters["max_area"]):
        return None
    if circularity < parameters["min_circularity"]:
        return None
    if solidity < parameters["min_solidity"]:
        return None
    if aspect_ratio > parameters["max_aspect_ratio"]:
        return None
    return {
        "contour": c,
        "area_px": area,
        "type": "vacuole",
        "circularity": circularity,
        "solidity": solidity,
        "aspect_ratio": aspect_ratio,
    }


def _contour_overlap(c1, c2, h, w):
    """Return overlap relative to the smaller contour area."""
    x1, y1, ww1, hh1 = cv2.boundingRect(c1)
    x2, y2, ww2, hh2 = cv2.boundingRect(c2)
    xa, ya = max(x1, x2), max(y1, y2)
    xb, yb = min(x1 + ww1, x2 + ww2), min(y1 + hh1, y2 + hh2)
    if xb <= xa or yb <= ya:
        return 0.0

    # To avoid allocating a full-image mask, compute the intersection only inside the bounding box.
    sub_w, sub_h = xb - xa, yb - ya
    m1 = np.zeros((sub_h, sub_w), np.uint8)
    m2 = np.zeros((sub_h, sub_w), np.uint8)
    c1_shift = c1.copy()
    c2_shift = c2.copy()
    c1_shift[:, :, 0] -= xa
    c1_shift[:, :, 1] -= ya
    c2_shift[:, :, 0] -= xa
    c2_shift[:, :, 1] -= ya
    cv2.drawContours(m1, [c1_shift], -1, 255, -1)
    cv2.drawContours(m2, [c2_shift], -1, 255, -1)
    inter = np.count_nonzero((m1 > 0) & (m2 > 0))
    a1 = np.count_nonzero(m1)
    a2 = np.count_nonzero(m2)
    return inter / max(min(a1, a2), 1)




def _safe_cell_interior_mask(cell_mask):
    """Returns the cell interior with a minimum margin equal to 2% of the cell scale."""
    area_cel = max(float(cv2.countNonZero(cell_mask)), 1.0)
    equivalent_radius = np.sqrt(area_cel / np.pi)
    margin = max(2.0, int(round(equivalent_radius * CONFIG.VACUOLE_MIN_BORDER_MARGIN_FRAC)))
    k = _odd_kernel(2 * margin + 1, 3)
    return cv2.erode(
        cell_mask,
        cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (k, k)),
    )


def _vacuole_shape_score(candidate, area_fraction):
    """Strict score based on circularity, solidity, and oval aspect ratio.

    Size helps break ties but never compensates for poor shape.
    """
    circ = candidate["circularity"]
    sol = candidate["solidity"]
    asp = candidate["aspect_ratio"]

    # Progressive penalty for highly elongated shapes.
    aspect_score = max(0.0, min(1.0, (3.8 - asp) / 2.8))
    # Continuous shape score: circularity, solidity, and ovality.
    shape_score = 0.45 * circ + 0.35 * sol + 0.20 * aspect_score
    size_score = min(1.0, area_fraction / 0.12)
    return 0.88 * shape_score + 0.12 * size_score




def _refine_outer_cell_contour(contour, shape):
    """Close inward green-boundary indentations using a fixed kernel near 50 px."""
    h, w = shape
    base = np.zeros((h, w), np.uint8)
    cv2.drawContours(base, [contour], -1, 255, -1)
    initial_area = max(cv2.contourArea(contour), 1.0)

    k0 = _odd_kernel(CONFIG.CELL_REENTRANCE_CLOSE_PX, 5)
    kernel_tests = []
    for k in (k0, max(5, k0 - 10), max(5, k0 - 20)):
        k = _odd_kernel(k, 5)
        if k not in kernel_tests:
            kernel_tests.append(k)

    for k in kernel_tests:
        closed = cv2.morphologyEx(
            base, cv2.MORPH_CLOSE,
            cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (k, k))
        )

        ks = _odd_kernel(CONFIG.CELL_EDGE_SMOOTH_PX, 3)
        if ks >= 3:
            closed = cv2.morphologyEx(
                closed, cv2.MORPH_CLOSE,
                cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (ks, ks))
            )

        cs, _ = cv2.findContours(closed, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_NONE)
        if not cs:
            continue

        M = cv2.moments(contour)
        if M["m00"] != 0:
            cx, cy = M["m10"]/M["m00"], M["m01"]/M["m00"]
            valid = [c for c in cs if cv2.pointPolygonTest(c, (cx, cy), False) >= 0]
            new_contour = max(valid, key=cv2.contourArea) if valid else max(cs, key=cv2.contourArea)
        else:
            new_contour = max(cs, key=cv2.contourArea)

        if cv2.contourArea(new_contour) <= initial_area * (1.0 + CONFIG.CELL_OUTER_MAX_EXPANSION):
            return new_contour

    return contour








def _classify_cell_gray_levels(base_image, cell_mask):
    """Automatically separate dark cytoplasm, gray vacuole, and bright lipid-like regions.

    Use 1D k-means on intensities ONLY inside the cell.
    Returns ordered centers and gray-interval limits.
    """
    display_values = base_image[cell_mask > 0].astype(np.float32)
    if display_values.size < 100:
        return None

    # Remove rare extremes so nearly black/white pixels do not dominate.
    p1, p99 = np.percentile(display_values, [1, 99])
    sample = display_values[(display_values >= p1) & (display_values <= p99)]
    if sample.size < 100:
        sample = display_values

    # Sample at most ~40k pixels to keep k-means fast on large TIFF files.
    if sample.size > 40000:
        idx = np.linspace(0, sample.size - 1, 40000).astype(np.int64)
        sample = np.sort(sample)[idx]

    data = sample.reshape(-1, 1).astype(np.float32)
    K = int(max(3, CONFIG.VACUOLE_TONE_K))
    criterios = (cv2.TERM_CRITERIA_EPS + cv2.TERM_CRITERIA_MAX_ITER, 80, 0.2)

    try:
        _, _, centers = cv2.kmeans(
            data, K, None, criterios, 8, cv2.KMEANS_PP_CENTERS
        )
        centers = sorted(float(x) for x in centers.ravel())
    except Exception:
        centers = []

    # The final classification uses three levels. If K > 3, condense them:
    # lower = cytoplasm, higher = lipid, median internal center = vacuole.
    if len(centers) >= 3:
        dark = centers[0]
        bright = centers[-1]
        internal_centers = centers[1:-1]
        mid = float(np.median(internal_centers)) if internal_centers else (dark + bright) / 2
    else:
        q25, q55, q90 = np.percentile(display_values, [25, 55, 90])
        dark, mid, bright = float(q25), float(q55), float(q90)

    # If the centers are too close, use quantiles for stabilization.
    if (mid - dark) < CONFIG.VACUOLE_TONE_MIN_GAP or (bright - mid) < CONFIG.VACUOLE_TONE_MIN_GAP:
        q20, q55, q92 = np.percentile(display_values, [20, 55, 92])
        dark = min(dark, float(q20))
        mid = float(q55)
        bright = max(bright, float(q92))

    # Class-midpoint limits exclude extreme dark and bright regions.
    lower = (dark + mid) / 2.0 - CONFIG.VACUOLE_GRAY_LOWER_PAD
    upper = (mid + bright) / 2.0 + CONFIG.VACUOLE_GRAY_UPPER_PAD

    # Guard against unreasonably wide intervals.
    p10, p95 = np.percentile(display_values, [10, 95])
    lower = max(lower, float(p10))
    upper = min(upper, float(p95))
    if upper <= lower + 4:
        lower = float(np.percentile(display_values, 35))
        upper = float(np.percentile(display_values, 82))

    return {
        "dark": dark,
        "mid": mid,
        "bright": bright,
        "lower": float(lower),
        "upper": float(upper),
    }











def _ellipse_contour_if_appropriate(c, shape, cell_mask):
    """Regularize only candidates that are already clearly oval.

    This removes small indentations/jagged threshold edges without turning
    every shape into an artificial ellipse.
    """
    if len(c) < 5:
        return c, 0.0

    area = max(cv2.contourArea(c), 1.0)
    hull = cv2.convexHull(c)
    solidity = area / max(cv2.contourArea(hull), 1.0)
    if solidity < CONFIG.VACUOLE_ELLIPSE_MIN_SOLIDITY:
        return c, 0.0

    try:
        el = cv2.fitEllipse(c)
    except Exception:
        return c, 0.0

    h, w = shape
    mc = np.zeros((h, w), np.uint8)
    cv2.drawContours(mc, [c], -1, 255, -1)

    me = np.zeros((h, w), np.uint8)
    cv2.ellipse(me, el, 255, -1)

    inter = cv2.countNonZero(cv2.bitwise_and(mc, me))
    union = cv2.countNonZero(cv2.bitwise_or(mc, me))
    iou = inter / max(union, 1)

    area_e = cv2.countNonZero(me)
    if iou < CONFIG.VACUOLE_MIN_ELLIPSE_IOU:
        return c, iou
    if area_e > area * (1.0 + CONFIG.VACUOLE_ELLIPSE_MAX_EXPANSION):
        return c, iou

    # The final ellipse must remain 100% inside the cell.
    outside_pixels = cv2.countNonZero(cv2.bitwise_and(me, cv2.bitwise_not(cell_mask)))
    if outside_pixels > 0:
        return c, iou

    ce, _ = cv2.findContours(me, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    if not ce:
        return c, iou
    ce = max(ce, key=cv2.contourArea)
    return ce, iou





def _detect_vacuoles_classical(base_image, enhanced, smoothed, th_val,
                       cell_mask, cell_area, cell_contour):
    """Single red-contour detector calibrated by TONE + TEXTURE inside the cell.

    Biological rules used:
      - BLACK/dark: cytoplasm/dense material -> reject
      - INTERMEDIATE gray: possible vacuole -> search
      - WHITE/very bright: lipid/inclusion -> reject
      - vacuole should be relatively homogeneous and round/oval
    """
    safe_mask = _safe_cell_interior_mask(cell_mask)
    display_values = base_image[cell_mask > 0].astype(np.float32)
    if display_values.size < 100:
        return []

    # ---------------------------------------------------------
    # 1) TONE CALIBRATION PER CELL
    # ---------------------------------------------------------
    # Instead of using a fixed gray value, thresholds adapt to each image/cell.
    p_low = float(np.percentile(display_values, CONFIG.VACUOLE_TONE_LOW_PERCENTILE))
    p_high = float(np.percentile(display_values, CONFIG.VACUOLE_TONE_HIGH_PERCENTILE))
    p_dark = float(np.percentile(display_values, 10))
    p_white = float(np.percentile(display_values, 92))

    # K-means is still calculated as a diagnostic/stabilizing step.
    tones = _classify_cell_gray_levels(base_image, cell_mask)
    if tones is not None:
        # Do not let the lower limit rise too far: darker gray vacuoles
        # remain eligible, while the upper limit is kept below very bright regions.
        lower = min(p_low, tones["lower"])
        upper = min(p_high, max(tones["upper"], float(np.percentile(display_values, 78))))
    else:
        lower, upper = p_low, p_high

    lower = max(lower, p_dark)
    upper = min(upper, p_white)
    if upper <= lower + 6:
        lower = float(np.percentile(display_values, 18))
        upper = float(np.percentile(display_values, 85))

    # ---------------------------------------------------------
    # 2) LOCAL TEXTURE
    # ---------------------------------------------------------
    # A gray vacuole tends to be smoother than granular/dense cytoplasm.
    x, y, bw, bh = cv2.boundingRect(cell_contour)
    diam = max(bw, bh)
    window = _odd_kernel(max(15, int(round(diam * CONFIG.VACUOLE_TEXTURE_WINDOW_FRAC))), 15)

    f = base_image.astype(np.float32)
    mean_img = cv2.boxFilter(f, -1, (window, window), normalize=True)
    mean_sq_img = cv2.boxFilter(f * f, -1, (window, window), normalize=True)
    local_std = np.sqrt(np.maximum(mean_sq_img - mean_img * mean_img, 0))

    vals_std = local_std[safe_mask > 0]
    if vals_std.size == 0:
        return []
    texture_threshold = float(np.percentile(vals_std, CONFIG.VACUOLE_TEXTURE_PERCENTILE))

    # Primary mask: GRAY TONE + LOW/MODERATE TEXTURE.
    support = (
        (base_image >= lower) &
        (base_image <= upper) &
        (local_std <= texture_threshold) &
        (safe_mask > 0)
    ).astype(np.uint8) * 255

    kclose = _odd_kernel(CONFIG.VACUOLE_GRAY_CLOSE_KERNEL, 3)
    kopen = _odd_kernel(CONFIG.VACUOLE_GRAY_OPEN_KERNEL, 3)
    mask = cv2.morphologyEx(
        support, cv2.MORPH_CLOSE,
        cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (kclose, kclose))
    )
    mask = cv2.morphologyEx(
        mask, cv2.MORPH_OPEN,
        cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (kopen, kopen))
    )
    mask = cv2.bitwise_and(mask, safe_mask)

    contours, _ = cv2.findContours(mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

    parameters = {
        "type": "vacuole",
        "min_area": CONFIG.MIN_VACUOLE_AREA_FRAC,
        "max_area": CONFIG.MAX_VACUOLE_AREA_FRAC,
        "min_circularity": CONFIG.MIN_VACUOLE_CIRCULARITY,
        "min_solidity": CONFIG.MIN_VACUOLE_SOLIDITY,
        "max_aspect_ratio": CONFIG.MAX_VACUOLE_ASPECT_RATIO,
        "min_score": CONFIG.MIN_VACUOLE_SHAPE_SCORE,
    }

    candidates = []
    for c0 in contours:
        # First size filter; the raw boundary may be irregular.
        initial_area = cv2.contourArea(c0)
        initial_area_fraction = initial_area / max(cell_area, 1e-9)
        if not (CONFIG.MIN_VACUOLE_AREA_FRAC <= initial_area_fraction <= CONFIG.MAX_VACUOLE_AREA_FRAC):
            continue

        # Regularize ONLY when the object already strongly matches an ellipse.
        c, ellipse_iou = _ellipse_contour_if_appropriate(
            c0, cell_mask.shape, cell_mask
        )

        candidate = _accept_vacuole_candidate(c, cell_area, parameters)
        if candidate is None:
            continue
        area_fraction = candidate["area_px"] / max(cell_area, 1e-9)

        # Fully inside the cell and, when not very large, inside the safe internal margin.
        mask_c = np.zeros_like(cell_mask)
        cv2.drawContours(mask_c, [c], -1, 255, -1)
        if cv2.countNonZero(cv2.bitwise_and(mask_c, cv2.bitwise_not(cell_mask))) > 0:
            continue
        if area_fraction < 0.30:
            if cv2.countNonZero(cv2.bitwise_and(mask_c, cv2.bitwise_not(safe_mask))) > 0:
                continue

        # Tone statistics inside the FINAL contour.
        vals = base_image[mask_c > 0].astype(np.float32)
        if vals.size < 30:
            continue
        med = float(np.median(vals))
        q10, q90 = np.percentile(vals, [10, 90])
        intensity_spread = float(q90 - q10)

        gray_support = float(np.mean((vals >= lower) & (vals <= upper)))
        dark_fraction = float(np.mean(vals < p_dark))
        white_fraction = float(np.mean(vals > p_white))

        # Very bright lipid-like regions and dark cytoplasm are explicitly excluded.
        min_support = (
            CONFIG.VACUOLE_GRAY_SMALL_MIN_SUPPORT
            if area_fraction < CONFIG.VACUOLE_STRICT_SMALL_BELOW
            else CONFIG.VACUOLE_GRAY_MIN_SUPPORT
        )
        max_spread = (
            CONFIG.VACUOLE_GRAY_SMALL_MAX_SPREAD
            if area_fraction < CONFIG.VACUOLE_STRICT_SMALL_BELOW
            else CONFIG.VACUOLE_GRAY_MAX_SPREAD
        )

        if gray_support < min_support:
            continue
        if white_fraction > CONFIG.VACUOLE_MAX_WHITE_FRACTION:
            continue
        if dark_fraction > CONFIG.VACUOLE_MAX_DARK_FRACTION:
            continue
        if intensity_spread > max_spread:
            continue
        if med <= p_dark or med >= p_white:
            continue

        # Small candidates still require particularly good shape.
        if area_fraction < CONFIG.VACUOLE_STRICT_SMALL_BELOW:
            if candidate["circularity"] < CONFIG.VACUOLE_SMALL_MIN_CIRCULARITY:
                continue
            if candidate["solidity"] < CONFIG.VACUOLE_SMALL_MIN_SOLIDITY:
                continue
            if candidate["aspect_ratio"] > CONFIG.VACUOLE_SMALL_MAX_ASPECT_RATIO:
                continue

        score = _vacuole_shape_score(candidate, area_fraction)
        score += 0.07 * gray_support
        score += 0.04 * min(1.0, ellipse_iou / max(CONFIG.VACUOLE_MIN_ELLIPSE_IOU, 1e-6))
        score += 0.03 * max(0.0, 1.0 - intensity_spread / max(max_spread, 1.0))

        candidate.update({
            "score": float(score),
            "source": "gray_texture",
            "gray_support": gray_support,
            "dark_frac": dark_fraction,
            "white_frac": white_fraction,
            "robust_spread": intensity_spread,
            "median_intensity": med,
            "ellipse_iou": float(ellipse_iou),
            "tone_lower": float(lower),
            "tone_upper": float(upper),
            "texture_threshold": float(texture_threshold),
        })

        if score >= parameters["min_score"]:
            candidates.append(candidate)

    candidates.sort(key=lambda v: (-v.get("score", 0), -v["area_px"]))

    # Deduplication.
    selected = []
    for cand in candidates:
        duplicate = False
        for other in selected:
            overlap = _contour_overlap(
                cand["contour"], other["contour"],
                cell_mask.shape[0], cell_mask.shape[1]
            )
            if overlap >= 0.45:
                duplicate = True
                break
        if not duplicate:
            selected.append(cand)
        if len(selected) >= CONFIG.MAX_VACUOLES_PER_CELL:
            break

    return selected






def _shape_local_brightness(c, image, cell_mask):
    H, W = image.shape
    mask = np.zeros((H, W), np.uint8)
    cv2.drawContours(mask, [c], -1, 255, -1)

    x, y, bw, bh = cv2.boundingRect(c)
    ring_px = max(3, int(round(max(bw, bh) * CONFIG.SHAPE_RING_FRAC)))
    ring_px = _odd_kernel(ring_px, 3)

    dil = cv2.dilate(
        mask,
        cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (ring_px, ring_px))
    )
    ring_mask = cv2.subtract(dil, mask)
    ring_mask = cv2.bitwise_and(ring_mask, cell_mask)

    vin = image[mask > 0]
    vr = image[ring_mask > 0]
    if vin.size < 30 or vr.size < 30:
        return None

    return {
        "delta_local": float(np.median(vin) - np.median(vr)),
        "median_in": float(np.median(vin)),
        "median_ring": float(np.median(vr)),
    }


def _detect_vacuoles_by_shape(base_image, cell_mask, cell_area):
    """Shape first; intensity only relative to the local cytoplasm."""
    vals = base_image[cell_mask > 0]
    if vals.size < 100:
        return []

    area_mask = max(float(cv2.countNonZero(cell_mask)), 1.0)
    radius = (area_mask / np.pi) ** 0.5
    margin = max(1, int(round(radius * 0.006)))
    ksafe = _odd_kernel(2*margin + 1, 3)
    safe = cv2.erode(
        cell_mask,
        cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (ksafe, ksafe))
    )

    candidates = []

    for perc in CONFIG.SHAPE_THRESHOLD_PERCENTILES:
        th = float(np.percentile(vals, perc))
        bw = ((base_image >= th) & (safe > 0)).astype(np.uint8) * 255

        kc = _odd_kernel(CONFIG.SHAPE_MORPH_CLOSE_PX, 3)
        ko = _odd_kernel(CONFIG.SHAPE_MORPH_OPEN_PX, 3)
        bw = cv2.morphologyEx(
            bw, cv2.MORPH_CLOSE,
            cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (kc, kc))
        )
        bw = cv2.morphologyEx(
            bw, cv2.MORPH_OPEN,
            cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (ko, ko))
        )

        cs, _ = cv2.findContours(bw, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_NONE)

        for c in cs:
            area, circ, sol, asp = _contour_metrics(c)
            area_fraction = area / max(cell_area, 1e-9)

            if not (CONFIG.SHAPE_MIN_AREA_FRAC <= area_fraction <= CONFIG.SHAPE_MAX_AREA_FRAC):
                continue
            if circ < CONFIG.SHAPE_MIN_CIRCULARITY:
                continue
            if sol < CONFIG.SHAPE_MIN_SOLIDITY:
                continue
            if asp > CONFIG.SHAPE_MAX_ASPECT_RATIO:
                continue

            if area_fraction < CONFIG.SHAPE_STRICT_BELOW_FRAC:
                if circ < CONFIG.SHAPE_SMALL_MIN_CIRCULARITY:
                    continue
                if sol < CONFIG.SHAPE_SMALL_MIN_SOLIDITY:
                    continue
                if asp > CONFIG.SHAPE_SMALL_MAX_ASPECT_RATIO:
                    continue

            mc = np.zeros_like(cell_mask)
            cv2.drawContours(mc, [c], -1, 255, -1)
            if cv2.countNonZero(cv2.bitwise_and(mc, cv2.bitwise_not(cell_mask))) > 0:
                continue

            rel = _shape_local_brightness(c, base_image, cell_mask)
            if rel is None or rel["delta_local"] < CONFIG.SHAPE_MIN_LOCAL_BRIGHTNESS_DELTA:
                continue

            score = (
                0.46 * min(1.0, max(0.0, circ)) +
                0.36 * min(1.0, max(0.0, sol)) +
                0.18 * min(1.0, 1.0/max(asp, 1.0))
            )
            candidates.append({
                "contour": c,
                "area_px": float(area),
                "type": "vacuole",
                "circularity": float(circ),
                "solidity": float(sol),
                "aspect_ratio": float(asp),
                "delta_local": rel["delta_local"],
                "score": float(score),
                "source": f"forma_p{perc}",
            })

    candidates.sort(key=lambda z: (z.get("score", 0), z.get("area_px", 0)), reverse=True)
    selected = []
    H, W = cell_mask.shape

    for cand in candidates:
        duplicado = False
        for other in selected:
            ov = _contour_overlap(cand["contour"], other["contour"], H, W)
            ratio = cand["area_px"] / max(other["area_px"], 1e-9)
            if ov >= 0.90 and 0.75 <= ratio <= 1.33:
                duplicado = True
                break
        if not duplicado:
            selected.append(cand)

    return selected

def _merge_hybrid_vacuoles(classical_list, shape_list, cell_mask):
    """Merge supported detectors; partial overlaps are not discarded."""
    all_candidates = []
    for candidates, source, bonus in (
        (classical_list, "CLASSICAL", 0.00),
        (shape_list, "SHAPE", 0.01),
    ):
        for v in (candidates or []):
            z = dict(v)
            z["hybrid_source"] = source + ":" + str(z.get("source", ""))
            z["_merge_score"] = float(z.get("score", 0.0)) + bonus
            all_candidates.append(z)

    all_candidates.sort(key=lambda z: (z.get("_merge_score", 0), z.get("area_px", 0)), reverse=True)

    selected = []
    H, W = cell_mask.shape

    for cand in all_candidates:
        duplicate = False

        for j, other in enumerate(selected):
            ov = _contour_overlap(cand["contour"], other["contour"], H, W)
            ratio = cand.get("area_px", 0) / max(other.get("area_px", 0), 1e-9)

            if (
                ov >= CONFIG.HYBRID_DUPLICATE_OVERLAP and
                CONFIG.HYBRID_DUPLICATE_AREA_RATIO_MIN <= ratio <=
                CONFIG.HYBRID_DUPLICATE_AREA_RATIO_MAX
            ):
                duplicate = True
                if cand.get("_merge_score", 0) > other.get("_merge_score", 0):
                    selected[j] = cand
                break

        if not duplicate:
            selected.append(cand)

    selected.sort(key=lambda z: (z.get("_merge_score", 0), z.get("area_px", 0)), reverse=True)
    return selected[:CONFIG.MAX_VACUOLES_PER_CELL]




def _detect_vacuoles(base_image, enhanced, smoothed, th_val,
                       cell_mask, cell_area, cell_contour):
    """Hybrid detector with a fast cascade.

    The classical and shape-based detectors run before their candidates are
    deduplicated and ranked.
    """
    classical_results, shape_results = [], []

    if CONFIG.HYBRID_USE_V24:
        try:
            classical_results = _detect_vacuoles_classical(
                base_image, enhanced, smoothed, th_val,
                cell_mask, cell_area, cell_contour
            )
        except Exception as e:
            print(f"[Hybrid] V24 detector failed for this cell: {e}")

    if CONFIG.SHAPE_DETECTOR_ENABLED:
        try:
            shape_results = _detect_vacuoles_by_shape(
                base_image, cell_mask, cell_area
            )
        except Exception as e:
            print(f"[Hybrid] Shape detector failed for this cell: {e}")

    return _merge_hybrid_vacuoles(classical_results, shape_results, cell_mask)




def _segment_vacuole_from_seed(image, cell_mask, cell_area, point):
    """Grow one rounded/oval vacuole from a user-selected internal point."""
    h, w = image.shape
    x, y = int(round(point[0])), int(round(point[1]))
    if not (0 <= x < w and 0 <= y < h) or cell_mask[y, x] == 0:
        return None

    smooth = cv2.GaussianBlur(image, (5, 5), 0)
    r = max(2, int(CONFIG.MANUAL_SEED_PATCH_RADIUS))
    y0, y1 = max(0, y-r), min(h, y+r+1)
    x0, x1 = max(0, x-r), min(w, x+r+1)
    patch_mask = cell_mask[y0:y1, x0:x1] > 0
    patch = smooth[y0:y1, x0:x1][patch_mask]
    if patch.size == 0:
        return None
    seed_tone = float(np.median(patch))

    cell_diameter = 2.0 * np.sqrt(max(cell_area, 1.0) / np.pi)
    k = _odd_kernel(max(3, round(CONFIG.MANUAL_SEED_SMOOTH_FRAC * cell_diameter)), 3)
    kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (k, k))
    candidates = []

    for tolerance in CONFIG.MANUAL_SEED_TOLERANCES:
        similar_mask = (
            (np.abs(smooth.astype(np.float32) - seed_tone) <= float(tolerance)) &
            (cell_mask > 0)
        ).astype(np.uint8) * 255
        similar_mask = cv2.morphologyEx(similar_mask, cv2.MORPH_CLOSE, kernel)
        similar_mask = cv2.morphologyEx(similar_mask, cv2.MORPH_OPEN, kernel)

        n, labels, stats, _ = cv2.connectedComponentsWithStats(similar_mask, 8)
        component_label = int(labels[y, x])
        if component_label <= 0 or component_label >= n:
            continue
        component_mask = (labels == component_label).astype(np.uint8) * 255
        contours, _ = cv2.findContours(component_mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        if not contours:
            continue
        contour = max(contours, key=cv2.contourArea)
        area = float(cv2.contourArea(contour))
        area_fraction = area / max(cell_area, 1.0)
        if not (CONFIG.MANUAL_SEED_MIN_AREA_FRAC <= area_fraction <= CONFIG.MANUAL_SEED_MAX_AREA_FRAC):
            continue

        perimeter = cv2.arcLength(contour, True)
        circularity = 4.0 * np.pi * area / max(perimeter * perimeter, 1e-9)
        hull_area = cv2.contourArea(cv2.convexHull(contour))
        solidity = area / max(hull_area, 1e-9)
        (_, _), (rw, rh), _ = cv2.minAreaRect(contour)
        aspect_ratio = max(rw, rh) / max(min(rw, rh), 1e-6)
        if (circularity < CONFIG.MANUAL_SEED_MIN_CIRCULARITY or
                solidity < CONFIG.MANUAL_SEED_MIN_SOLIDITY or
                aspect_ratio > CONFIG.MANUAL_SEED_MAX_ASPECT_RATIO):
            continue

        # Prefer a stable, large, rounded region; do not force a perfect circle.
        score = (0.50 * min(circularity, 1.0) + 0.35 * min(solidity, 1.0) +
                 0.15 * min(area_fraction / 0.20, 1.0))
        candidates.append({
            "contour": contour,
            "area_px": area,
            "score": score,
            "source": "manual_seed",
            "seed": (x, y),
            "tolerance": float(tolerance),
        })

    if not candidates:
        return None
    candidates.sort(key=lambda c: (c["score"], c["area_px"]), reverse=True)
    return candidates[0]



def _segment_cell_from_seed(image, seed):
    """Recover a complete cell contour from one click near the cell center."""
    h, w = image.shape
    x, y = map(int, seed)
    if not (0 <= x < w and 0 <= y < h):
        return None

    clahe = cv2.createCLAHE(
        clipLimit=CONFIG.CLAHE_CLIP,
        tileGridSize=(CONFIG.CLAHE_TILE, CONFIG.CLAHE_TILE),
    )
    enhanced = clahe.apply(image)
    blur_k = _odd_kernel(CONFIG.GAUSS_BLUR_CELL, 3)
    smoothed = cv2.GaussianBlur(enhanced, (blur_k, blur_k), 0)
    otsu_value, _ = cv2.threshold(
        smoothed, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU
    )

    local_radius = max(25, int(round(0.04 * min(h, w))))
    x0, x1 = max(0, x-local_radius), min(w, x+local_radius+1)
    y0, y1 = max(0, y-local_radius), min(h, y+local_radius+1)
    local = smoothed[y0:y1, x0:x1]
    local_median = float(np.median(local)) if local.size else float(smoothed[y, x])

    thresholds = sorted(set(int(round(v)) for v in [
        otsu_value,
        np.percentile(smoothed, 35),
        np.percentile(smoothed, 45),
        np.percentile(smoothed, 55),
        min(245, local_median + 8),
        min(245, local_median + 18),
    ]))

    candidates = []
    image_area = float(h*w)

    for threshold in thresholds:
        dark = (smoothed <= threshold).astype(np.uint8) * 255

        min_component = max(20, int(round(CONFIG.DENOISE_MIN_COMPONENT_FRAC * image_area)))
        n_comp, labels, stats, _ = cv2.connectedComponentsWithStats(dark, connectivity=8)
        cleaned = np.zeros_like(dark)
        for label_id in range(1, n_comp):
            if stats[label_id, cv2.CC_STAT_AREA] >= min_component:
                cleaned[labels == label_id] = 255

        close_size = _odd_kernel(max(15, min(CONFIG.CLOSE_KERNEL_CELL, 61)), 3)
        open_size = _odd_kernel(max(3, min(CONFIG.OPEN_KERNEL_CELL, 9)), 3)
        close_kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (close_size, close_size))
        open_kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (open_size, open_size))
        cleaned = cv2.morphologyEx(cleaned, cv2.MORPH_CLOSE, close_kernel)
        cleaned = cv2.morphologyEx(cleaned, cv2.MORPH_OPEN, open_kernel)

        contours, _ = cv2.findContours(cleaned, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        for contour in contours:
            area = float(cv2.contourArea(contour))
            area_fraction = area / image_area
            if not (0.002 <= area_fraction <= 0.95):
                continue
            if cv2.pointPolygonTest(contour, (float(x), float(y)), False) < 0:
                continue

            hull_area = float(cv2.contourArea(cv2.convexHull(contour)))
            solidity = area / max(hull_area, 1e-9)
            (_, _), (rw, rh), _ = cv2.minAreaRect(contour)
            aspect = max(rw, rh) / max(min(rw, rh), 1e-6)
            if solidity < 0.40 or aspect > 5.0:
                continue

            refined = _refine_outer_cell_contour(contour, (h, w))
            refined_area = float(cv2.contourArea(refined))
            if refined_area <= 0:
                continue

            score = (
                0.55 * min(solidity, 1.0)
                + 0.25 * (1.0 / max(aspect, 1.0))
                + 0.20 * (1.0 - min(refined_area / image_area, 0.95))
            )
            candidates.append((score, refined_area, refined))

    if not candidates:
        return None
    candidates.sort(key=lambda item: (-item[0], item[1]))
    return candidates[0][2]


def _review_cells_interactively(image, initial_cells, image_name=""):
    """Mandatory cell-review step shown before vacuole review for every image."""
    selected = []
    for cell in initial_cells:
        contour = cell["cell_contour"].copy()
        selected.append({
            "cell_contour": contour,
            "cell_area_px": float(cv2.contourArea(contour)),
            "source": cell.get("source", "automatic"),
        })

    state = {
        "mode": "A",
        "drawing": False,
        "points": [],
        "draft_artist": None,
        "split_lines": [],
    }
    fig, ax = plt.subplots(figsize=(11, 9))
    ax.imshow(image, cmap="gray")
    ax.set_axis_off()
    artists = []

    def redraw(message=""):
        nonlocal artists
        for artist in artists:
            try:
                artist.remove()
            except Exception:
                pass
        artists = []

        for idx, cell in enumerate(selected, start=1):
            c = cell["cell_contour"][:, 0, :]
            closed = np.vstack([c, c[0]])
            line, = ax.plot(closed[:, 0], closed[:, 1], color="lime", linewidth=2.0)
            artists.append(line)
            m = cv2.moments(cell["cell_contour"])
            if m["m00"]:
                txt = ax.text(
                    m["m10"]/m["m00"], m["m01"]/m["m00"], f"C{idx}",
                    color="lime", fontsize=11, weight="bold"
                )
                artists.append(txt)

        # Pending split lines are shown in yellow until they are applied.
        for split_idx, split_line in enumerate(state["split_lines"], start=1):
            if len(split_line) >= 2:
                xy = np.asarray(split_line, dtype=float)
                line, = ax.plot(
                    xy[:, 0], xy[:, 1],
                    color="yellow", linewidth=2.2, linestyle="--"
                )
                artists.append(line)
                txt = ax.text(
                    xy[len(xy)//2, 0], xy[len(xy)//2, 1],
                    f"D{split_idx}", color="yellow", fontsize=10, weight="bold"
                )
                artists.append(txt)

        if state["mode"] == "A":
            instructions = "CELL MODE A — CLICK near the CENTER of a missing cell | B: draw | C: erase | D: split"
        elif state["mode"] == "B":
            instructions = "CELL MODE B — HOLD LEFT BUTTON and DRAW the green cell boundary | A: seed | C: erase | D: split"
        elif state["mode"] == "C":
            instructions = "CELL MODE C — CLICK inside a green contour to ERASE | A: seed | B: draw | D: split"
        else:
            instructions = (
                "CELL MODE D — DRAW one or MORE yellow cut lines across merged cells. "
                "Press ENTER to APPLY cuts | A/B/C to leave split mode"
            )

        status = message or f"Selected: {len(selected)} cell(s)"
        if not selected:
            status += " — At least ONE cell is required before ENTER."

        ax.set_title(
            f"{image_name}\n{instructions}\n"
            "RIGHT CLICK/DELETE: undo last | S: clear all | ENTER: approve (or apply D cuts)\n"
            + status
        )
        fig.canvas.draw_idle()

    def cancel_draft():
        state["drawing"] = False
        state["points"] = []
        if state["draft_artist"] is not None:
            try:
                state["draft_artist"].remove()
            except Exception:
                pass
            state["draft_artist"] = None

    def add_contour(contour, source):
        if contour is None or len(contour) < 3:
            return False, "No valid cell contour was produced."

        area = float(cv2.contourArea(contour))
        image_area = float(image.shape[0] * image.shape[1])

        # Manual drawing is intentionally permissive: if the user explicitly
        # drew a cell, do not reject it using the stricter automatic thresholds.
        if source == "manual_draw":
            min_area = max(25.0, 0.0002 * image_area)
            max_area = 0.98 * image_area
        else:
            min_area = 0.0015 * image_area
            max_area = 0.95 * image_area

        if not (min_area <= area <= max_area):
            return False, (
                f"The cell was drawn, but its area ({area:.0f} px²) is outside "
                "the permitted range. Try drawing around the full cell."
            )

        # Strong duplicate protection remains, but manual drawing is allowed to
        # overlap moderately with an automatic contour so the user can correct it.
        duplicate_limit = 0.90 if source == "manual_draw" else 0.70
        for existing in selected:
            overlap = _contour_overlap(
                contour, existing["cell_contour"], image.shape[0], image.shape[1]
            )
            if overlap > duplicate_limit:
                return False, (
                    "This drawn cell almost completely overlaps an existing green cell. "
                    "Erase the incorrect one with C first, then draw again with B."
                )

        # Manual review may reveal more cells than the automatic default.
        # Do not silently reject a manually drawn real cell.
        if len(selected) >= CONFIG.MAX_CELLS_PER_IMAGE and source != "manual_draw":
            return False, f"Maximum cells per image reached ({CONFIG.MAX_CELLS_PER_IMAGE})."

        selected.append({
            "cell_contour": contour,
            "cell_area_px": area,
            "source": source,
        })
        return True, None

    def contour_from_drawing(points):
        """Convert a free-hand stroke into a CLOSED green cell contour.

        The user does not need to return exactly to the starting point.
        The program automatically connects the last point back to the first.
        """
        if len(points) < 3:
            return None, "Draw around the cell before releasing the mouse."

        pts = np.asarray(points, dtype=np.float32)
        pts = pts[np.r_[True, np.any(np.diff(pts, axis=0) != 0, axis=1)]]
        if len(pts) < 3:
            return None, "The drawn contour is too short."

        pts_i = np.round(pts).astype(np.int32)

        # Explicitly close the contour: last point -> first point.
        if not np.array_equal(pts_i[0], pts_i[-1]):
            pts_i = np.vstack([pts_i, pts_i[0]])

        polygon = np.zeros(image.shape, np.uint8)

        # Draw the closed outline first, then fill it. This is more robust than
        # relying only on fillPoly for irregular free-hand strokes.
        outline_width = max(2, int(round(0.0015 * min(image.shape))))
        cv2.polylines(
            polygon,
            [pts_i.reshape(-1, 1, 2)],
            isClosed=True,
            color=255,
            thickness=outline_width,
        )
        cv2.fillPoly(
            polygon,
            [pts_i.reshape(-1, 1, 2)],
            255,
        )

        # Fill any tiny gaps created by a rough mouse stroke.
        area = float(np.count_nonzero(polygon))
        if area <= 0:
            return None, "The drawn contour has no enclosed area."

        diameter = 2.0 * np.sqrt(area / np.pi)
        # Keep smoothing light so the manual contour follows what the user drew.
        k = _odd_kernel(max(3, min(11, round(0.004 * diameter))), 3)
        kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (k, k))
        polygon = cv2.morphologyEx(polygon, cv2.MORPH_CLOSE, kernel)

        contours, _ = cv2.findContours(
            polygon, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE
        )
        if not contours:
            return None, "The drawn boundary could not be converted into a cell."

        contour = max(contours, key=cv2.contourArea)
        if cv2.contourArea(contour) <= 0:
            return None, "The drawn boundary has zero area."

        return contour, None

    def _draw_split_line_on_mask(mask, points):
        """Cut a binary cell mask along a user-drawn polyline."""
        if len(points) < 2:
            return mask
        pts = np.round(np.asarray(points, dtype=np.float32)).astype(np.int32)
        # A small but visible cut width; enough to separate touching contours
        # without deleting a large part of either cell.
        cut_width = max(3, int(round(0.0025 * min(image.shape))))
        if cut_width % 2 == 0:
            cut_width += 1
        out = mask.copy()
        cv2.polylines(out, [pts.reshape(-1, 1, 2)], False, 0, thickness=cut_width)
        return out

    def apply_split_lines():
        """Apply all pending D-lines, allowing 2, 3, 4... merged cells to split."""
        if not state["split_lines"]:
            return False, "No split lines have been drawn."

        changed = False
        new_selected = []

        for cell in selected:
            original_contour = cell["cell_contour"]
            cell_mask = np.zeros(image.shape, dtype=np.uint8)
            cv2.drawContours(cell_mask, [original_contour], -1, 255, -1)

            # Only use cut lines that actually intersect this cell.
            relevant = []
            for line in state["split_lines"]:
                pts = np.round(np.asarray(line, dtype=np.float32)).astype(np.int32)
                if len(pts) < 2:
                    continue
                hit = False
                for px, py in pts:
                    if 0 <= px < image.shape[1] and 0 <= py < image.shape[0]:
                        if cell_mask[py, px] > 0:
                            hit = True
                            break
                if hit:
                    relevant.append(line)

            if not relevant:
                new_selected.append(cell)
                continue

            cut_mask = cell_mask.copy()
            for line in relevant:
                cut_mask = _draw_split_line_on_mask(cut_mask, line)

            # A cut must fully separate components. Connected components is more
            # reliable here than contour hierarchy for manually cut masks.
            n_labels, labels, stats, _ = cv2.connectedComponentsWithStats(
                (cut_mask > 0).astype(np.uint8), connectivity=8
            )

            pieces = []
            original_area = max(float(cv2.contourArea(original_contour)), 1.0)
            min_piece_area = max(
                40.0,
                0.025 * original_area,
                0.0008 * float(image.shape[0] * image.shape[1]),
            )

            for label_id in range(1, n_labels):
                area_px = float(stats[label_id, cv2.CC_STAT_AREA])
                if area_px < min_piece_area:
                    continue

                piece_mask = np.zeros_like(cell_mask)
                piece_mask[labels == label_id] = 255

                # Restore only a tiny amount of edge removed by the cutting line.
                # This smooths the new boundary without reconnecting the pieces.
                k = 3
                kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (k, k))
                piece_mask = cv2.morphologyEx(piece_mask, cv2.MORPH_CLOSE, kernel)

                contours_piece, _ = cv2.findContours(
                    piece_mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE
                )
                if not contours_piece:
                    continue
                contour_piece = max(contours_piece, key=cv2.contourArea)
                piece_area = float(cv2.contourArea(contour_piece))
                if piece_area < min_piece_area:
                    continue

                pieces.append({
                    "cell_contour": contour_piece,
                    "cell_area_px": piece_area,
                    "source": "manual_split",
                })

            if len(pieces) >= 2:
                changed = True
                new_selected.extend(pieces)
            else:
                # Keep the original if the line did not fully separate it.
                new_selected.append(cell)

        if changed:
            if len(new_selected) > CONFIG.MAX_CELLS_PER_IMAGE:
                # Preserve all split results: the user explicitly indicated them.
                # Raise the runtime limit for this image instead of silently
                # discarding a manually created cell.
                pass
            selected[:] = new_selected
            state["split_lines"].clear()
            return True, f"Split applied. There are now {len(selected)} cell(s)."

        return False, (
            "The cut line(s) did not completely separate a merged cell. "
            "Draw each yellow line fully across the green contour and press ENTER again."
        )

    def undo_last():
        if state["drawing"]:
            cancel_draft()
            redraw("Current drawing cancelled.")
            return
        if state["mode"] == "D" and state["split_lines"]:
            state["split_lines"].pop()
            redraw("Last split line removed.")
            return
        if selected:
            selected.pop()
            redraw("Last cell removed.")

    def on_click(event):
        if event.inaxes is not ax:
            return
        if event.button == 3:
            undo_last()
            return
        if event.button != 1 or event.xdata is None or event.ydata is None:
            return

        x, y = float(event.xdata), float(event.ydata)

        if state["mode"] == "C":
            matches = []
            for index, cell in enumerate(selected):
                if cv2.pointPolygonTest(cell["cell_contour"], (x, y), False) >= 0:
                    matches.append((float(cv2.contourArea(cell["cell_contour"])), index))
            if not matches:
                redraw("No green cell contour was found at this point.")
                return
            _, index = min(matches)
            selected.pop(index)
            redraw("Cell contour erased.")
            return

        if state["mode"] in ("B", "D"):
            cancel_draft()
            state["drawing"] = True
            state["points"] = [(x, y)]
            state["draft_artist"], = ax.plot(
                [x], [y], color="yellow",
                linewidth=2.2 if state["mode"] == "D" else 1.6,
                linestyle="--" if state["mode"] == "D" else "-"
            )
            if state["mode"] == "D":
                redraw(
                    "Drawing split line... cross the entire merged green cell and release. "
                    "You may draw additional D-lines before pressing ENTER."
                )
            else:
                redraw("Drawing cell... trace around the cell and release. The contour will close automatically.")
            return

        for cell in selected:
            if cv2.pointPolygonTest(cell["cell_contour"], (x, y), False) >= 0:
                redraw("This point is already inside a selected green cell.")
                return

        contour = _segment_cell_from_seed(image, (int(round(x)), int(round(y))))
        if contour is None:
            redraw("Automatic recovery failed. Click closer to the center or press B and draw the cell.")
            return

        ok, error = add_contour(contour, "seed")
        redraw("Cell added from click." if ok else error)

    def on_motion(event):
        if not state["drawing"] or event.inaxes is not ax:
            return
        if event.xdata is None or event.ydata is None:
            return
        point = (float(event.xdata), float(event.ydata))
        if state["points"]:
            px, py = state["points"][-1]
            if (point[0]-px)**2 + (point[1]-py)**2 < 1.5**2:
                return
        state["points"].append(point)
        xy = np.asarray(state["points"])
        state["draft_artist"].set_data(xy[:, 0], xy[:, 1])
        fig.canvas.draw_idle()

    def on_release(event):
        if not state["drawing"]:
            return
        # Some Matplotlib backends report button=None when the mouse is released
        # just outside the image axes. If a drawing is active, finalize it anyway.
        if event.button not in (1, None):
            return

        points = list(state["points"])
        current_mode = state["mode"]
        cancel_draft()

        if current_mode == "D":
            if len(points) < 2:
                redraw("Split line too short. Draw it fully across the merged cell.")
                return
            state["split_lines"].append(points)
            redraw(
                f"Split line {len(state['split_lines'])} stored. "
                "Draw another line for additional merged cells, or press ENTER to apply all cuts."
            )
            return

        contour, error = contour_from_drawing(points)
        if contour is None:
            redraw(error)
            return
        ok, error = add_contour(contour, "manual_draw")
        redraw("Drawn cell contour added." if ok else error)

    def on_key(event):
        if event.key in ("enter", "return"):
            if state["mode"] == "D" and state["split_lines"]:
                cancel_draft()
                ok, message = apply_split_lines()
                redraw(message)
                # ENTER in D mode applies the cuts but intentionally keeps the
                # cell-review window open so the user can inspect C1/C2/C3...
                return
            if not selected:
                redraw("At least one cell is required. Use A or B before continuing.")
                return
            cancel_draft()
            plt.close(fig)
        elif event.key in ("backspace", "delete"):
            undo_last()
        elif event.key and event.key.lower() == "a":
            cancel_draft()
            state["mode"] = "A"
            state["split_lines"].clear()
            redraw("Mode A activated: click near the center of a missing cell.")
        elif event.key and event.key.lower() == "b":
            cancel_draft()
            state["mode"] = "B"
            state["split_lines"].clear()
            redraw("Mode B activated: draw the green cell boundary.")
        elif event.key and event.key.lower() == "c":
            cancel_draft()
            state["mode"] = "C"
            state["split_lines"].clear()
            redraw("Mode C activated: click inside a green contour to erase it.")
        elif event.key and event.key.lower() == "d":
            cancel_draft()
            state["mode"] = "D"
            redraw(
                "Mode D activated: draw one or more yellow cut lines. "
                "For 3 merged cells, draw 2 cuts, then press ENTER."
            )
        elif event.key and event.key.lower() == "s":
            cancel_draft()
            state["split_lines"].clear()
            selected.clear()
            redraw("All cell contours removed.")

    fig.canvas.mpl_connect("button_press_event", on_click)
    fig.canvas.mpl_connect("motion_notify_event", on_motion)
    fig.canvas.mpl_connect("button_release_event", on_release)
    fig.canvas.mpl_connect("key_press_event", on_key)

    redraw()
    plt.tight_layout()
    plt.show(block=True)
    return selected


def _review_vacuoles_interactively(image, cells, image_name=""):
    """Review automatic contours; A adds by seed, B draws, and C erases."""
    if not cells:
        return {}
    masks = []
    for cell in cells:
        mask = np.zeros(image.shape, np.uint8)
        cv2.drawContours(mask, [cell["cell_contour"]], -1, 255, -1)
        masks.append(mask)

    selected = []  # (cell_index, candidate, matplotlib_artist)
    state = {"mode": "A", "drawing": False, "points": [], "draft_artist": None}
    fig, ax = plt.subplots(figsize=(11, 9))
    ax.imshow(image, cmap="gray")
    for idx, cell in enumerate(cells, 1):
        c = cell["cell_contour"][:, 0, :]
        ax.plot(c[:, 0], c[:, 1], color="lime", linewidth=1.4)
        m = cv2.moments(cell["cell_contour"])
        if m["m00"]:
            ax.text(m["m10"]/m["m00"], m["m01"]/m["m00"], f"C{idx}",
                    color="lime", fontsize=11, weight="bold")

    for cell_index, cell in enumerate(cells):
        for candidate in cell.get("vacuoles", []):
            c = candidate["contour"][:, 0, :]
            closed = np.vstack([c, c[0]])
            artist, = ax.plot(closed[:, 0], closed[:, 1], color="red", linewidth=2.0)
            selected.append((cell_index, candidate, artist))

    base_title = f"{image_name}"
    ax.set_axis_off()

    def redraw_title(message=""):
        if state["mode"] == "A":
            instructions = "MODE A — CLICK inside a missing vacuole | B: draw | C: erase"
        elif state["mode"] == "B":
            instructions = "MODE B — HOLD LEFT BUTTON and DRAW | A: seed | C: erase"
        else:
            instructions = "MODE C — CLICK inside a red contour to ERASE | A: seed | B: draw"
        status = message or f"Selected: {len(selected)} vacuole(s)"
        ax.set_title(
            base_title + "\n" + instructions +
            "\nRIGHT CLICK/DELETE: remove last | ENTER: approve image | S: remove all" +
            "\n" + status
        )
        fig.canvas.draw_idle()

    def cancel_draft():
        state["drawing"] = False
        state["points"] = []
        if state["draft_artist"] is not None:
            try:
                state["draft_artist"].remove()
            except Exception:
                pass
            state["draft_artist"] = None

    def candidate_from_drawing(points):
        if len(points) < 12:
            return None, None, "Draw a complete contour while holding the left button."
        pts = np.asarray(points, dtype=np.float32)
        pts = pts[np.r_[True, np.any(np.diff(pts, axis=0) != 0, axis=1)]]
        if len(pts) < 6:
            return None, None, "The drawn contour is too short."
        x0, y0 = np.mean(pts, axis=0)
        cell_index = next((i for i, mask in enumerate(masks)
                           if 0 <= int(y0) < mask.shape[0] and 0 <= int(x0) < mask.shape[1]
                           and mask[int(y0), int(x0)] > 0), None)
        if cell_index is None:
            return None, None, "Draw the contour inside one green cell."
        if sum(1 for i, _, _ in selected if i == cell_index) >= CONFIG.MAX_VACUOLES_PER_CELL:
            return None, None, f"Cell C{cell_index+1} already has the maximum number of vacuoles."

        polygon = np.zeros(image.shape, np.uint8)
        cv2.fillPoly(polygon, [np.round(pts).astype(np.int32).reshape(-1, 1, 2)], 255)
        polygon = cv2.bitwise_and(polygon, masks[cell_index])
        diameter = 2.0 * np.sqrt(max(cells[cell_index]["cell_area_px"], 1.0) / np.pi)
        k = _odd_kernel(max(3, round(0.009 * diameter)), 3)
        kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (k, k))
        polygon = cv2.morphologyEx(polygon, cv2.MORPH_CLOSE, kernel)
        polygon = cv2.morphologyEx(polygon, cv2.MORPH_OPEN, kernel)
        contours, _ = cv2.findContours(polygon, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        if not contours:
            return None, None, "The drawn boundary could not be closed. Try again."
        contour = max(contours, key=cv2.contourArea)
        area = float(cv2.contourArea(contour))
        area_cell = max(cells[cell_index]["cell_area_px"], 1.0)
        area_fraction = area / area_cell
        if not (CONFIG.MANUAL_SEED_MIN_AREA_FRAC <= area_fraction <= CONFIG.MANUAL_SEED_MAX_AREA_FRAC):
            return None, None, "The drawn area is outside the permitted vacuole size range."
        perimeter = cv2.arcLength(contour, True)
        circularity = 4.0 * np.pi * area / max(perimeter * perimeter, 1e-9)
        hull_area = cv2.contourArea(cv2.convexHull(contour))
        solidity = area / max(hull_area, 1e-9)
        candidate = {
            "contour": contour, "area_px": area,
            "score": 0.5 * min(circularity, 1.0) + 0.5 * min(solidity, 1.0),
            "source": "manual_draw", "seed": None, "tolerance": None,
        }
        return cell_index, candidate, None

    def undo_last():
        if state["drawing"]:
            cancel_draft()
            redraw_title("Current drawing cancelled.")
            return
        if selected:
            _, _, artist = selected.pop()
            artist.remove()
            redraw_title("Last selection removed.")

    def on_click(event):
        if event.inaxes is not ax:
            return
        if event.button == 3:
            undo_last()
            return
        if event.button != 1 or event.xdata is None or event.ydata is None:
            return
        if state["mode"] == "C":
            x, y = float(event.xdata), float(event.ydata)
            matches = []
            for index, (_, candidate, _) in enumerate(selected):
                contour = candidate["contour"]
                if cv2.pointPolygonTest(contour, (x, y), False) >= 0:
                    matches.append((float(cv2.contourArea(contour)), index))
            if not matches:
                redraw_title("No red contour was found at this point.")
                return
            _, index = min(matches)
            _, _, artist = selected.pop(index)
            artist.remove()
            redraw_title("Contour erased.")
            return
        if state["mode"] == "B":
            cancel_draft()
            state["drawing"] = True
            state["points"] = [(float(event.xdata), float(event.ydata))]
            state["draft_artist"], = ax.plot(
                [event.xdata], [event.ydata], color="yellow", linewidth=1.6
            )
            redraw_title("Drawing... release the left button to close the contour.")
            return
        x, y = int(round(event.xdata)), int(round(event.ydata))
        cell_index = next((i for i, mask in enumerate(masks)
                           if 0 <= y < mask.shape[0] and 0 <= x < mask.shape[1]
                           and mask[y, x] > 0), None)
        if cell_index is None:
            redraw_title("The click is outside every green cell. Try again.")
            return
        if sum(1 for i, _, _ in selected if i == cell_index) >= CONFIG.MAX_VACUOLES_PER_CELL:
            redraw_title(f"Cell C{cell_index+1} already has the maximum number of vacuoles.")
            return
        candidate = _segment_vacuole_from_seed(
            image, masks[cell_index], cells[cell_index]["cell_area_px"], (x, y)
        )
        if candidate is None:
            redraw_title("No rounded boundary was found from this point. Click nearer the center.")
            return
        for old_cell, old, _ in selected:
            if old_cell == cell_index:
                overlap = _contour_overlap(
                    candidate["contour"], old["contour"], image.shape[0], image.shape[1]
                )
                if overlap > 0.65:
                    redraw_title("This vacuole is already selected.")
                    return
        c = candidate["contour"][:, 0, :]
        artist, = ax.plot(c[:, 0], c[:, 1], color="red", linewidth=2.0)
        selected.append((cell_index, candidate, artist))
        redraw_title()

    def on_motion(event):
        if not state["drawing"] or event.inaxes is not ax:
            return
        if event.xdata is None or event.ydata is None:
            return
        point = (float(event.xdata), float(event.ydata))
        if state["points"]:
            px, py = state["points"][-1]
            if (point[0]-px)**2 + (point[1]-py)**2 < 1.5**2:
                return
        state["points"].append(point)
        xy = np.asarray(state["points"])
        state["draft_artist"].set_data(xy[:, 0], xy[:, 1])
        fig.canvas.draw_idle()

    def on_release(event):
        if not state["drawing"] or event.button != 1:
            return
        points = list(state["points"])
        cancel_draft()
        cell_index, candidate, error = candidate_from_drawing(points)
        if candidate is None:
            redraw_title(error)
            return
        c = candidate["contour"][:, 0, :]
        closed = np.vstack([c, c[0]])
        artist, = ax.plot(closed[:, 0], closed[:, 1], color="red", linewidth=2.0)
        selected.append((cell_index, candidate, artist))
        redraw_title("Drawn contour added.")

    def on_key(event):
        if event.key in ("enter", "return"):
            cancel_draft()
            plt.close(fig)
        elif event.key in ("backspace", "delete"):
            undo_last()
        elif event.key and event.key.lower() == "a":
            cancel_draft()
            state["mode"] = "A"
            redraw_title("Mode A activated: click inside a vacuole.")
        elif event.key and event.key.lower() == "b":
            cancel_draft()
            state["mode"] = "B"
            redraw_title("Mode B activated: hold the left button and draw the contour.")
        elif event.key and event.key.lower() == "c":
            cancel_draft()
            state["mode"] = "C"
            redraw_title("Mode C activated: click inside a red contour to erase it.")
        elif event.key and event.key.lower() == "s":
            cancel_draft()
            selected.clear()
            plt.close(fig)

    fig.canvas.mpl_connect("button_press_event", on_click)
    fig.canvas.mpl_connect("motion_notify_event", on_motion)
    fig.canvas.mpl_connect("button_release_event", on_release)
    fig.canvas.mpl_connect("key_press_event", on_key)
    redraw_title()
    plt.tight_layout()
    plt.show(block=True)

    result = {i: [] for i in range(len(cells))}
    for cell_index, candidate, _ in selected:
        result[cell_index].append(candidate)
    return result


def segment_cells_and_vacuoles(image, image_name=""):
    """Automatic proposals + mandatory cell review + mandatory vacuole review."""
    h, w = image.shape
    image_area = h * w

    clahe = cv2.createCLAHE(
        clipLimit=CONFIG.CLAHE_CLIP,
        tileGridSize=(CONFIG.CLAHE_TILE, CONFIG.CLAHE_TILE),
    )
    enhanced = clahe.apply(image)
    smoothed = cv2.GaussianBlur(
        enhanced, (CONFIG.GAUSS_BLUR_CELL, CONFIG.GAUSS_BLUR_CELL), 0
    )
    th_val, threshold_mask = cv2.threshold(
        smoothed, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU
    )

    raw_mask = (threshold_mask == 0).astype(np.uint8) * 255

    minimum_component_area = CONFIG.DENOISE_MIN_COMPONENT_FRAC * image_area
    n_comp, labels, stats, _ = cv2.connectedComponentsWithStats(raw_mask, connectivity=8)
    denoised_mask = np.zeros_like(raw_mask)
    for i in range(1, n_comp):
        if stats[i, cv2.CC_STAT_AREA] >= minimum_component_area:
            denoised_mask[labels == i] = 255

    k_close = cv2.getStructuringElement(
        cv2.MORPH_ELLIPSE, (_odd_kernel(CONFIG.CLOSE_KERNEL_CELL, 3),) * 2
    )
    mask = cv2.morphologyEx(denoised_mask, cv2.MORPH_CLOSE, k_close)

    k_open = cv2.getStructuringElement(
        cv2.MORPH_ELLIPSE, (_odd_kernel(CONFIG.OPEN_KERNEL_CELL, 3),) * 2
    )
    mask = cv2.morphologyEx(mask, cv2.MORPH_OPEN, k_open)

    k_fill = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (_odd_kernel(31, 3),) * 2)
    mask = cv2.morphologyEx(mask, cv2.MORPH_CLOSE, k_fill)

    temporary_contours, _ = cv2.findContours(mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    filled_mask = np.zeros_like(mask)
    for contour in temporary_contours:
        area = cv2.contourArea(contour)
        if CONFIG.MIN_CELL_AREA_FRAC * image_area <= area <= CONFIG.MAX_CELL_AREA_FRAC * image_area:
            cv2.drawContours(filled_mask, [contour], -1, 255, -1)
    mask = filled_mask

    contours, hierarchy = cv2.findContours(mask, cv2.RETR_TREE, cv2.CHAIN_APPROX_SIMPLE)

    automatic_cells = []

    if hierarchy is not None:
        hierarchy = hierarchy[0]
        cell_candidates = []

        for i, contour in enumerate(contours):
            if hierarchy[i][3] != -1:
                continue
            area = cv2.contourArea(contour)
            area_fraction = area / image_area
            if not (CONFIG.MIN_CELL_AREA_FRAC <= area_fraction <= CONFIG.MAX_CELL_AREA_FRAC):
                continue

            hull = cv2.convexHull(contour)
            hull_area = cv2.contourArea(hull)
            solidity = area / hull_area if hull_area > 0 else 0
            if solidity < CONFIG.MIN_CELL_SOLIDITY:
                continue

            (_, _), (rw, rh), _ = cv2.minAreaRect(contour)
            aspect_ratio = max(rw, rh) / max(min(rw, rh), 1e-6)
            if aspect_ratio > CONFIG.MAX_CELL_ASPECT_RATIO:
                continue

            cell_candidates.append((area, i))

        cell_candidates.sort(key=lambda item: -item[0])
        accepted = []

        for area, i in cell_candidates:
            moments = cv2.moments(contours[i])
            if moments["m00"] == 0:
                continue
            cx = moments["m10"] / moments["m00"]
            cy = moments["m01"] / moments["m00"]
            if any(
                cv2.pointPolygonTest(contours[j], (cx, cy), False) >= 0
                for _, j in accepted
            ):
                continue
            accepted.append((area, i))
            if len(accepted) >= CONFIG.MAX_CELLS_PER_IMAGE:
                break

        for _, idx in accepted:
            cell_contour = _refine_outer_cell_contour(contours[idx], (h, w))
            automatic_cells.append({
                "cell_contour": cell_contour,
                "cell_area_px": float(cv2.contourArea(cell_contour)),
                "vacuoles": [],
                "source": "automatic",
            })

    # NEW: this window ALWAYS opens, even when automatic_cells is empty.
    reviewed_cells = _review_cells_interactively(image, automatic_cells, image_name)

    if not reviewed_cells:
        raise RuntimeError(
            "No cell was approved. At least one cell must be selected before vacuole analysis."
        )

    result = []
    for reviewed_cell in reviewed_cells:
        cell_contour = reviewed_cell["cell_contour"]
        cell_area = float(cv2.contourArea(cell_contour))
        cell_mask = np.zeros((h, w), dtype=np.uint8)
        cv2.drawContours(cell_mask, [cell_contour], -1, 255, -1)

        detected_vacuoles = _detect_vacuoles(
            image, enhanced, smoothed, th_val,
            cell_mask, cell_area, cell_contour
        )
        detected_vacuoles = list(detected_vacuoles)
        detected_vacuoles.sort(key=lambda v: (-v.get("score", 0), -v["area_px"]))
        detected_vacuoles = detected_vacuoles[:CONFIG.MAX_VACUOLES_PER_CELL]

        result.append({
            "cell_contour": cell_contour,
            "cell_area_px": cell_area,
            "vacuoles": detected_vacuoles,
            "cell_source": reviewed_cell.get("source", "reviewed"),
        })

    selected_vacuoles = _review_vacuoles_interactively(image, result, image_name)
    for i, cell in enumerate(result):
        cell["vacuoles"] = selected_vacuoles.get(i, [])[:CONFIG.MAX_VACUOLES_PER_CELL]

    return result


# =================================================================
# 4) VERIFICATION IMAGE DRAWING (green=cell, red=vacuole)
# =================================================================
def draw_annotations(image, cells):
    """Verification image: green = cell; red = vacuole."""
    annotated = cv2.cvtColor(image, cv2.COLOR_GRAY2BGR)
    line_thickness = max(2, image.shape[0] // 500)
    for i, cell in enumerate(cells, start=1):
        cv2.drawContours(annotated, [cell["cell_contour"]], -1, (0, 255, 0), line_thickness)
        M = cv2.moments(cell["cell_contour"])
        if M["m00"] != 0:
            cx, cy = int(M["m10"] / M["m00"]), int(M["m01"] / M["m00"])
        else:
            cx, cy = cell["cell_contour"][0][0]
        cv2.putText(annotated, f"C{i}", (cx - 20, cy),
                    cv2.FONT_HERSHEY_SIMPLEX, 1.4, (0, 255, 0), line_thickness)
        for v in cell["vacuoles"]:
            cv2.drawContours(annotated, [v["contour"]], -1, (0, 0, 255), line_thickness)

    fs = max(0.55, image.shape[0] / 1800)
    y = max(25, line_thickness * 8)
    cv2.putText(annotated, "RED = vacuole", (15, y),
                cv2.FONT_HERSHEY_SIMPLEX, fs, (255, 255, 255),
                max(1, line_thickness // 2), cv2.LINE_AA)
    return annotated





# =================================================================
# 5) PARAMETER WINDOW (edit without changing the code)
# =================================================================
# Each item: (CONFIG_attribute_name, label, explanation, type)
CELL_PARAMETER_FIELDS = [
    ("CELL_REENTRANCE_CLOSE_PX", "Close cell indentations (px)",
     "Default: 51 px. Closes deep inward notches in the green cell boundary. Increase it if the contour enters too far into the cell; decrease it if nearby cells start to merge.", int),
    ("CELL_OUTER_MAX_EXPANSION", "Maximum cell expansion",
     "Limits how much the indentation-closing step may enlarge the cell mask. Lower values preserve the original outline more strictly.", float),
    ("CLOSE_KERNEL_CELL", "Cell closing kernel (px)",
     "Higher values join more distant pieces of the SAME cell (useful for interrupted cell masks). "
     "If too large, it may merge two nearby cells.", int),
    ("OPEN_KERNEL_CELL", "Cell opening kernel (px)",
     "Higher values remove thin bridges/tendrils connecting the cell to the background. "
     "If too large, it may remove legitimate thin cell regions.", int),
    ("MIN_CELL_AREA_FRAC", "Minimum cell area",
     "Fraction of image area (ex.: 0.01 = 1%). Lower values accept smaller cells.", float),
    ("MAX_CELL_AREA_FRAC", "Maximum cell area",
     "Fraction of image area. Higher values accept larger cells.", float),
    ("MIN_CELL_SOLIDITY", "Minimum cell solidity",
     "0 a 1. Higher values require a rounder/more filled cell (rejecting line-like background artifacts).", float),
    ("MAX_CELL_ASPECT_RATIO", "Maximum cell aspect ratio",
     "Lower values require rounder cells; higher values accept more elongated/oval cells.", float),
    ("MAX_CELLS_PER_IMAGE", "Maximum cells per image",
     "Maximum number of cells (largest first) accepted per image.", int),
    ("CELL_REENTRANCE_CLOSE_PX", "Close outer-boundary indentations",
     "Controls how strongly inward notches are closed so the green boundary does not enter bright vacuoles. "
     "Higher values close deeper indentations; no convex hull is used.", float),
    ("CELL_OUTER_MAX_EXPANSION", "Maximum cell expansion",
     "Limits how much external refinement can enlarge the cell area.", float),
]

VACUOLE_PARAMETER_FIELDS = [
    ("MANUAL_SEED_PATCH_RADIUS", "Seed sampling radius (px)",
     "Radius used to measure the vacuole gray tone around your click. Increase slightly for noisy images; decrease for very small vacuoles.", int),
    ("MANUAL_SEED_MIN_AREA_FRAC", "Minimum vacuole area",
     "Minimum fraction of the cell occupied by a clicked vacuole. Default 0.004 = 0.4%.", float),
    ("MANUAL_SEED_MAX_AREA_FRAC", "Maximum vacuole area",
     "Maximum fraction of the cell occupied by one vacuole. Default 0.90 = 90%.", float),
    ("MANUAL_SEED_MIN_CIRCULARITY", "Minimum circularity",
     "Increase to require a rounder boundary; decrease to accept more naturally deformed vacuoles.", float),
    ("MANUAL_SEED_MIN_SOLIDITY", "Minimum solidity",
     "Increase to reject contours with deep indentations or fragmented edges.", float),
    ("MANUAL_SEED_MAX_ASPECT_RATIO", "Maximum aspect ratio",
     "Maximum elongation of a vacuole. Lower values require a rounder or less elongated oval.", float),
    ("MANUAL_SEED_SMOOTH_FRAC", "Boundary smoothing",
     "Morphological smoothing relative to cell diameter. Increase for noisy borders; decrease to retain finer boundary detail.", float),
    ("MAX_VACUOLES_PER_CELL", "Maximum vacuoles per cell",
     "Maximum number of user-selected vacuoles accepted inside each detected cell.", int),
]



SCALE_BRIGHTNESS_PARAMETER_FIELDS = [
    ("SCALEBAR_LENGTH_PX", "Fallback scale-bar length (px)",
     "Used only as a fallback if automatic bar measurement fails; the bar is typically about 450-600 px in this dataset. "
     "Adjust if you know the expected bar length in pixels.", int),
    ("SCALEBAR_SEARCH_WIDTH_FRAC", "Scale-label search width",
     "Fraction of IMAGE WIDTH searched from the lower-left corner (e.g., 0.45 = "
     "45% of the width). INCREASE if the scale label is outside this "
     "region.", float),
    ("SCALEBAR_SEARCH_HEIGHT_FRAC", "Scale-label search height",
     "Fraction of IMAGE HEIGHT searched from the lower-left corner. INCREASE if the "
     "scale label is outside this region.", float),
    ("NORMALIZE_BRIGHTNESS", "Normalize brightness across images (1=yes, 0=no)",
     "If enabled, each image is normalized to a common grayscale range before "
     "segmentation; useful when a batch contains both brighter and darker images.", bool),
    ("BRIGHTNESS_LOW_PERCENTILE", "Low brightness percentile",
     "LOWER = normalization is more sensitive to isolated very dark pixels.", float),
    ("BRIGHTNESS_HIGH_PERCENTILE", "High brightness percentile",
     "HIGHER = normalization is more sensitive to isolated very bright pixels.", float),
    ("AUTO_GAMMA", "Automatic gamma for dark images (1=yes)",
     "If enabled, automatically brightens only images that remain dark after normalization.", bool),
    ("AUTO_GAMMA_TRIGGER_MEDIAN", "Gamma trigger median",
     "Gamma is applied only if the normalized image median is below this value.", float),
    ("AUTO_GAMMA_TARGET_MEDIAN", "Gamma target median",
     "Target median brightness when automatic gamma correction is needed.", float),
]


def edit_parameters_gui():
    """Open a tabbed window so the user can review/edit the
    main CONFIG parameters without editing the source code. Each
    field includes an explanation of what increasing/decreasing the value does.
    Update CONFIG.* directly when the user clicks 'Apply and continue'."""
    # Use a child Toplevel window attached to the existing main Tk window.
    # Creating a second tk.Tk() here can corrupt Tkinter state.
    # of Tkinter and can prevent subsequent dialogs from appearing
    # This can prevent subsequent dialogs from appearing.
    if tk._default_root is None:
        window = tk.Tk()  # Standalone use of this function when main() is not already running.
    else:
        window = tk.Toplevel()
    window.title("Segmentation parameters - review/adjust")
    window.geometry("760x560")

    aviso = tk.Label(
        window,
        text="Adjust the values if needed, or leave them unchanged, and click "
        "'Apply and continue'. Each field includes a short explanation.",
        wraplength=740, justify="left", font=("", 9, "italic"),
    )
    aviso.pack(padx=10, pady=(10, 5), anchor="w")

    notebook = ttk.Notebook(window)
    notebook.pack(fill="both", expand=True, padx=10, pady=5)

    entries = {}

    def build_tab(fields, title):
        sheet = ttk.Frame(notebook)
        notebook.add(sheet, text=title)

        canvas = tk.Canvas(sheet, borderwidth=0)
        scrollbar = ttk.Scrollbar(sheet, orient="vertical", command=canvas.yview)
        inner_frame = ttk.Frame(canvas)
        inner_frame.bind(
            "<Configure>", lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
        canvas.create_window((0, 0), window=inner_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")

        for r, (name, label_text, description, candidate_type) in enumerate(fields):
            current_value = getattr(CONFIG, name)
            tk.Label(inner_frame, text=label_text, font=("", 9, "bold"), anchor="w").grid(
                row=r, column=0, sticky="w", padx=5, pady=4
            )
            entry = tk.Entry(inner_frame, width=10)
            entry.insert(0, str(current_value))
            entry.grid(row=r, column=1, padx=5, pady=4)
            tk.Label(
                inner_frame, text=description, wraplength=460, justify="left",
                font=("", 8), fg="#444444",
            ).grid(row=r, column=2, sticky="w", padx=5, pady=4)
            entries[name] = (entry, candidate_type)

    build_tab(CELL_PARAMETER_FIELDS, "Cell")
    build_tab(VACUOLE_PARAMETER_FIELDS, "Vacuole")
    build_tab(SCALE_BRIGHTNESS_PARAMETER_FIELDS, "Scale / Brightness")

    result = {"ok": False}

    def apply_changes():
        try:
            for name, (entry, candidate_type) in entries.items():
                text = entry.get().strip()
                if candidate_type is bool:
                    value = text not in ("0", "False", "false", "")
                else:
                    value = candidate_type(text)
                setattr(CONFIG, name, value)
        except ValueError as e:
            messagebox.showerror("Invalid value", f"Check the entered values.\n{e}")
            return
        result["ok"] = True
        window.destroy()

    def cancel_changes():
        result["ok"] = False
        window.destroy()

    button_frame = tk.Frame(window)
    button_frame.pack(fill="x", padx=10, pady=10)
    tk.Button(button_frame, text="Apply and continue", command=apply_changes, bg="#4CAF50", fg="white").pack(
        side="right", padx=5
    )
    tk.Button(button_frame, text="Use default values (skip)", command=cancel_changes).pack(side="right")

    window.grab_set()  # make the window modal
    window.wait_window()
    return result["ok"]


# =================================================================
# 6) PROCESS ONE IMAGE (returns data rows + annotated image)
# =================================================================
def process_image(path, annotated_output_folder, calibration_state=None):
    name = Path(path).name
    image = load_grayscale_image(path)

    if calibration_state is None:
        calibration_state = {"mode": "ask"}

    # Try automatic calibration first for EVERY image.
    calibration = automatic_scale_calibration(image) if TESSERACT_OK else None

    if calibration is None:
        mode = calibration_state.get("mode", "ask")

        # Ask only ONCE for the batch, on the first failure.
        if mode == "ask":
            use_pixels = messagebox.askyesno(
                "Automatic calibration failed",
                "The scale could not be read automatically in this image:\\n\\n"
                f"{name}\\n\\n"
                "YES = process THIS image and ALL subsequent OCR failures in PIXELS, "
                "without asking again.\\n\\n"
                "NO = attempt MANUAL calibration for each image when OCR fails."
            )

            if use_pixels:
                calibration_state["mode"] = "pixels"
                mode = "pixels"
            else:
                calibration_state["mode"] = "manual"
                mode = "manual"

        if mode == "pixels":
            calibration = {
                "um_per_px": None,
                "method": "pixels",
                "details": "OCR failed; the batch is configured to process failures in pixels",
            }

        elif mode == "manual":
            calibration = manual_scale_calibration(image, name)

            # If manual calibration fails or is cancelled, offer to process the
            # remaining images in pixels.
            if calibration is None:
                switch_to_pixels = messagebox.askyesno(
                    "Manual calibration not completed",
                    "Manual calibration for this image was not completed.\\n\\n"
                    "Process THIS image and ALL subsequent failures in PIXELS?\\n\\n"
                    "YES = use pixels for the rest of the batch.\\n"
                    "NO = skip only this image."
                )

                if switch_to_pixels:
                    calibration_state["mode"] = "pixels"
                    calibration = {
                        "um_per_px": None,
                        "method": "pixels",
                        "details": "manual calibration failed; remaining OCR failures will be processed in pixels",
                    }
                else:
                    raise RuntimeError(
                        "Manual calibration cancelled; image skipped by the user."
                    )

    normalized_image = normalize_image_brightness(image)
    cells = segment_cells_and_vacuoles(normalized_image, name)
    annotated_image = draw_annotations(image, cells)

    # Save as high-quality JPEG (smaller than PNG while preserving
    # sufficient sharpness for visual verification and PowerPoint export).
    output_image_path = str(Path(annotated_output_folder) / f"{Path(name).stem}_anotada.jpg")
    cv2.imwrite(output_image_path, annotated_image, [cv2.IMWRITE_JPEG_QUALITY, 92])

    um_per_px = calibration["um_per_px"]
    rows = []
    for i, cell in enumerate(cells, start=1):
        cell_area_px = cell["cell_area_px"]
        vacuole_area_px = sum(v["area_px"] for v in cell["vacuoles"])
        num_vacuoles = len(cell["vacuoles"])
        ratio_pct = (vacuole_area_px / cell_area_px * 100) if cell_area_px > 0 else 0

        log_line = {
            "file": name,
            "cell": i,
            "num_vacuoles": num_vacuoles,
            "cell_area_px2": cell_area_px,
            "vacuole_area_px2": vacuole_area_px,
            "vacuole_cell_ratio_%": round(ratio_pct, 2),
            "calibration_um_per_px": um_per_px,
            "calibration_method": calibration["method"],
            "calibration_details": calibration["details"],
        }
        if um_per_px:
            log_line["cell_area_um2"] = round(cell_area_px * (um_per_px ** 2), 4)
            log_line["vacuole_area_um2"] = round(vacuole_area_px * (um_per_px ** 2), 4)
        else:
            log_line["cell_area_um2"] = None
            log_line["vacuole_area_um2"] = None

        rows.append(log_line)

    if not rows:
        rows.append(
            {
                "file": name, "cell": 0, "num_vacuoles": 0,
                "cell_area_px2": None, "vacuole_area_px2": None,
                "vacuole_cell_ratio_%": None,
                "calibration_um_per_px": um_per_px,
                "calibration_method": calibration["method"],
                "calibration_details": "NO CELL DETECTED - review image",
                "cell_area_um2": None, "vacuole_area_um2": None,
            }
        )

    return rows, output_image_path


# =================================================================
# 7) EXPORT: EXCEL + POWERPOINT (2 images per slide, side by side)
# =================================================================
def export_excel(all_rows, xlsx_path):
    dataframe = pd.DataFrame(all_rows)
    columns = [
        "file", "cell", "num_vacuoles",
        "cell_area_px2", "vacuole_area_px2", "vacuole_cell_ratio_%",
        "cell_area_um2", "vacuole_area_um2",
        "calibration_um_per_px", "calibration_method", "calibration_details",
    ]
    dataframe = dataframe[columns]
    with pd.ExcelWriter(xlsx_path, engine="openpyxl") as writer:
        dataframe.to_excel(writer, sheet_name="data_per_cell", index=False)

        summary = (
            dataframe.groupby("file")
            .agg(
                n_cells=("cell", "nunique"),
                total_vacuoles=("num_vacuoles", "sum"),
                mean_vacuole_ratio_pct=("vacuole_cell_ratio_%", "mean"),
            )
            .reset_index()
        )
        summary.to_excel(writer, sheet_name="summary_per_image", index=False)

        for sheet in writer.sheets.values():
            for column in sheet.columns:
                max_length = max((len(str(c.value)) for c in column if c.value is not None), default=10)
                sheet.column_dimensions[column[0].column_letter].width = min(max_length + 2, 40)


def _fill_slide_slot(slide, file_name, image_path, rows, x_offset):
    """Draw one block (title + image + table) in one half of the slide,
    comecando na posicao horizontal x_offset (Inches)."""
    block_width = Inches(6.35)

    title = slide.shapes.add_textbox(x_offset, Inches(0.15), block_width, Inches(0.4))
    title.text_frame.text = file_name
    title.text_frame.paragraphs[0].font.size = Pt(14)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.word_wrap = True

    image_height = Inches(4.55)
    picture = slide.shapes.add_picture(image_path, x_offset, Inches(0.62), height=image_height)
    if picture.width > block_width:
        slide.shapes._spTree.remove(picture._element)
        picture = slide.shapes.add_picture(image_path, x_offset, Inches(0.62), width=block_width)
    picture.left = int(x_offset + (block_width - picture.width) / 2)

    table_y = Inches(0.62) + image_height + Inches(0.1)
    num_rows = len(rows) + 1
    num_columns = 4
    row_height = Inches(0.3)
    table_shape = slide.shapes.add_table(
        num_rows, num_columns, x_offset, table_y, block_width, row_height * num_rows
    )
    table = table_shape.table
    headers = ["Cell", "No. vacuoles", "Cell area (µm²)", "Vacuole/Cell (%)"]
    for c, text in enumerate(headers):
        table.cell(0, c).text = text
        table.cell(0, c).text_frame.paragraphs[0].font.size = Pt(10)
        table.cell(0, c).text_frame.paragraphs[0].font.bold = True

    for r, log_line in enumerate(rows, start=1):
        area_um2 = log_line.get("cell_area_um2")
        display_values = [
            str(log_line["cell"]),
            str(log_line["num_vacuoles"]),
            f"{area_um2:.3f}" if area_um2 is not None else "N/D (px)",
            str(log_line["vacuole_cell_ratio_%"]),
        ]
        for c, text in enumerate(display_values):
            table.cell(r, c).text = text
            table.cell(r, c).text_frame.paragraphs[0].font.size = Pt(10)


def export_powerpoint(results_list, pptx_path):
    """results_list: list of tuples (file_name, annotated_image_path, data_rows).
    Generate 2 analyses per slide"""
    presentation = Presentation()
    presentation.slide_width = Inches(13.33)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]

    slide = presentation.slides.add_slide(blank_layout)
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(2))
    tf = textbox.text_frame
    tf.text = "Vacuole Analysis - C. vulgaris"
    tf.paragraphs[0].font.size = Pt(40)
    p2 = tf.add_paragraph()
    p2.text = "Green = cell | Red = vacuole"
    p2.font.size = Pt(20)

    slot_offsets = [Inches(0.15), Inches(6.83)]
    for pos in range(0, len(results_list), 2):
        pair = results_list[pos : pos + 2]
        slide = presentation.slides.add_slide(blank_layout)
        for slot, (file_name, image_path, rows) in enumerate(pair):
            _fill_slide_slot(slide, file_name, image_path, rows, slot_offsets[slot])

    presentation.save(pptx_path)


class ProgressWindow:
    """Simple window showing batch-processing progress -
    so the user can always see that the script is running (rather than
    appearing to disappear in a console window)."""

    def __init__(self, total):
        self.window = tk.Toplevel()
        self.window.title("Processing images...")
        self.window.geometry("520x170")
        self.window.protocol("WM_DELETE_WINDOW", lambda: None)  # prevent accidental closing
        self.window.attributes("-topmost", True)

        self.status_label = tk.Label(
            self.window, text="Starting...", font=("", 10), wraplength=490, justify="left"
        )
        self.status_label.pack(pady=(15, 5), padx=15, anchor="w")

        self.progress_bar = ttk.Progressbar(self.window, length=490, maximum=max(total, 1))
        self.progress_bar.pack(pady=5, padx=15)

        self.count_label = tk.Label(self.window, text=f"0 / {total}", font=("", 9))
        self.count_label.pack()

        self.error_label = tk.Label(self.window, text="", fg="#b00000", font=("", 9), wraplength=490)
        self.error_label.pack(pady=5, padx=15)

        self.total = total
        self.error_count = 0
        self.window.update()

    def update(self, idx, file_name, error_message=False):
        self.progress_bar["value"] = idx
        self.count_label.config(text=f"{idx} / {self.total}")
        if error_message:
            self.error_count += 1
            self.status_label.config(text=f"ERROR processing: {file_name}", fg="#b00000")
        else:
            self.status_label.config(text=f"Processing: {file_name}", fg="black")
        if self.error_count:
            self.error_label.config(
                text=f"{self.error_count} image(s) with errors so far (details in "
                f"processing_log.txt and the final summary)."
            )
        self.window.update()

    def close(self):
        self.window.destroy()


def _write_log(log_file, text):
    """Write a timestamped line to the console and the log file.

    Flush the log immediately so progress remains recorded even if the script
    crashes or its window closes.
    """
    from datetime import datetime

    log_line = f"[{datetime.now().strftime('%H:%M:%S')}] {text}"
    print(log_line)
    if log_file is not None:
        log_file.write(log_line + "\n")
        log_file.flush()


# =================================================================
# 8) MAIN PROGRAM
# =================================================================
def main():
    # Startup diagnostic: if Tkinter has a problem, show
    # a clear terminal message instead of closing silently.
    try:
        root = tk.Tk()
    except Exception as e:
        print("ERROR starting the Tkinter interface:", e)
        print("On Windows, reinstall Python with Tcl/Tk enabled and run: pip install matplotlib pillow")
        raise
    root.withdraw()

    input_folder = filedialog.askdirectory(title="Select the folder containing the images")
    if not input_folder:
        print("No folder selected. Exiting.")
        return

    output_folder = Path(input_folder) / "vacuole_results"
    annotated_output_folder = output_folder / "annotated_images"
    annotated_output_folder.mkdir(parents=True, exist_ok=True)

    extensions = (".tif", ".tiff", ".png", ".jpg", ".jpeg")
    files = sorted(
        [p for p in Path(input_folder).iterdir() if p.suffix.lower() in extensions]
    )

    if not files:
        messagebox.showwarning("Warning", "No supported image files were found in this folder.")
        return

    if not TESSERACT_OK:
        messagebox.showwarning(
            "OCR warning",
            "The pytesseract library was not found, or "
            "Tesseract-OCR is not installed on this computer.\n\n"
            "The script will still work, but it may ask you to ENTER "
            "the scale value manually when automatic "
            "OCR is unavailable.",
        )

    # --- 1) parameter window (review/adjust without editing the code) ---
    edit_parameters_gui()

    # --- 1.5) initial scale check ---
    if TESSERACT_OK:
        first_image = load_grayscale_image(str(files[0]))
        calibration_test = automatic_scale_calibration(first_image)
        if calibration_test is not None:
            messagebox.showinfo(
                "Scale recognized",
                "Automatic scale reading worked for the first image:\n\n"
                + calibration_test["details"]
                + "\n\nIf another image fails, manual calibration will be available."
            )
        else:
            messagebox.showwarning(
                "Scale OCR",
                "Automatic scale reading failed for the first image.\n\n"
                "The program will continue trying every image; each failure "
                "can be calibrated manually or processed in pixels depending on the selected batch mode."
            )
    else:
        messagebox.showwarning(
            "OCR unavailable",
            "Tesseract is unavailable. Images will require "
            "manual calibration or pixel-only processing."
        )

    # --- 2) mandatory GREEN cell review followed by RED vacuole review ---
    # Cell detection remains automatic. For each image, the user clicks once
    # inside every vacuole and the program finds its rounded boundary from the
    # local gray-tone difference, restricted to the detected cell.

    all_rows = []
    powerpoint_results = []
    errors = []

    total = len(files)

    log_path = output_folder / "processing_log.txt"
    log_file = open(log_path, "w", encoding="utf-8")
    _write_log(log_file, f"Starting processing of {total} image(s) from folder: {input_folder}")
    _write_log(log_file, f"Results will be saved to: {output_folder}")

    progress = ProgressWindow(total)
    # Calibration-failure handling mode shared across the entire batch.
    calibration_state = {"mode": "ask"}

    for idx, path in enumerate(files, start=1):
        progress.update(idx - 1, path.name)
        _write_log(log_file, f"[{idx}/{total}] Processing {path.name} ...")
        try:
            rows, annotated_image_path = process_image(str(path), annotated_output_folder, calibration_state)
            all_rows.extend(rows)
            powerpoint_results.append((path.name, annotated_image_path, rows))
            num_cells = len({l["cell"] for l in rows if l["cell"] != 0})
            _write_log(log_file, f"    OK - {num_cells} cell(s) found")
            progress.update(idx, path.name)
        except Exception as e:
            _write_log(log_file, f"    ERROR processing {path.name}: {e}")
            traceback.print_exc()
            log_file.write(traceback.format_exc() + "\n")
            log_file.flush()
            errors.append((path.name, str(e)))
            progress.update(idx, path.name, error_message=True)

    progress.close()

    if not all_rows:
        _write_log(log_file, "NO image was processed successfully. Exiting without creating Excel/PPTX files.")
        log_file.close()
        messagebox.showerror(
            "Error",
            f"No image was processed successfully.\n\n"
            f"See details in:\n{log_path}",
        )
        return

    _write_log(log_file, "Creating Excel spreadsheet...")
    xlsx_path = output_folder / "vacuole_results.xlsx"
    export_excel(all_rows, str(xlsx_path))

    _write_log(log_file, "Creating PowerPoint presentation...")
    pptx_path = output_folder / "vacuole_report.pptx"
    export_powerpoint(powerpoint_results, str(pptx_path))

    _write_log(log_file, "Processing completed.")
    log_file.close()

    message_text = (
        f"Processing completed!\n\n"
        f"Images processed: {total - len(errors)}/{total}\n"
        f"Excel: {xlsx_path}\n"
        f"PowerPoint: {pptx_path}\n"
        f"Annotated images: {annotated_output_folder}\n"
        f"Complete log: {log_path}"
    )
    if errors:
        message_text += "\n\nImages with errors:\n" + "\n".join(f"- {n}: {e}" for n, e in errors)

    print(message_text)
    messagebox.showinfo("Completed", message_text)


if __name__ == "__main__":
    try:
        main()
    except Exception:
        # If something unexpected happens outside the image-processing loop
        # (ex.: error opening the folder), show the error in a dialog instead
        # of closing without explanation.
        full_error = traceback.format_exc()
        print(full_error)
        try:
            messagebox.showerror(
                "Unexpected error",
                "The script stopped because of an unexpected error:\n\n" + full_error[-1500:],
            )
        except Exception:
            pass
    finally:
        # Keep the console window open when launched by double-click
        # in the file so any message can be read before closing.
        try:
            input("\nPress ENTER to close this window...")
        except Exception:
            pass
